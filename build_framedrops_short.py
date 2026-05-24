"""
Framedrops — 5-Slide Portrait (9:16) Reel/Short
Target: under 30 seconds total. ~5-6s per slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------- Palette ----------
NAVY        = RGBColor(0x08, 0x1B, 0x3A)
NAVY_DEEP   = RGBColor(0x05, 0x10, 0x26)
ROYAL       = RGBColor(0x00, 0x57, 0xFF)
PURPLE      = RGBColor(0x7B, 0x2C, 0xFF)
PINK        = RGBColor(0xFF, 0x4D, 0x8D)
ORANGE      = RGBColor(0xFF, 0x8A, 0x3D)
SOFT_WHITE  = RGBColor(0xF8, 0xFA, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0xB6, 0xC2, 0xE0)
CARD_FILL   = RGBColor(0x10, 0x24, 0x4A)
DIVIDER     = RGBColor(0x1B, 0x2E, 0x5C)

# ---------- Setup: 9:16 Portrait ----------
prs = Presentation()
prs.slide_width  = Inches(7.5)
prs.slide_height = Inches(13.333)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill_rgb, line=False, line_rgb=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb or WHITE
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_gradient_rect(slide, x, y, w, h, stops, angle=2700000, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    grad.set('flip', 'none'); grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle)); lin.set('scaled', '0')
    etree.SubElement(grad, qn('a:tileRect'))
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_gradient_line_rect(slide, x, y, w, h, stops, line_w_pt=1.5, angle=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.background()
    sp = s._element.spPr
    for el in sp.findall(qn('a:ln')):
        sp.remove(el)
    ln = etree.SubElement(sp, qn('a:ln'))
    ln.set('w', str(int(line_w_pt * 12700))); ln.set('cap', 'rnd')
    grad = etree.SubElement(ln, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle)); lin.set('scaled', '0')
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font='Helvetica Neue', italic=False, line_spacing=1.15, letter_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
        if letter_sp is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(letter_sp))
    return tb


def add_gradient_text(slide, x, y, w, h, text, size=44, bold=True, stops=None, align=PP_ALIGN.LEFT,
                      anchor=MSO_ANCHOR.TOP, font='Helvetica Neue', angle=0, line_spacing=1.05, letter_sp=None):
    if stops is None:
        stops = [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold
        rPr = r._r.get_or_add_rPr()
        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        grad = etree.SubElement(rPr, qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        for pos, rgb in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(int(pos * 100000)))
            srgb = etree.SubElement(gs, qn('a:srgbClr'))
            srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(angle)); lin.set('scaled', '0')
        if letter_sp is not None:
            rPr.set('spc', str(letter_sp))
    return tb


def add_blob(slide, cx, cy, r, rgb, alpha=40000):
    x = cx - r // 2; y = cy - r // 2
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, r, r)
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    sf = etree.SubElement(sp, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    alphaEl = etree.SubElement(srgb, qn('a:alpha'))
    alphaEl.set('val', str(alpha))
    s.line.fill.background(); s.shadow.inherit = False
    return s


def background(slide):
    add_rect(slide, 0, 0, SW, SH, NAVY_DEEP)
    add_gradient_rect(slide, 0, 0, SW, SH, stops=[(0.0, NAVY_DEEP), (1.0, NAVY)], angle=5400000)
    add_blob(slide, Inches(1.0), Inches(1.2), Inches(4.5), PURPLE, alpha=30000)
    add_blob(slide, Inches(6.8), Inches(3.5), Inches(4.0), ROYAL, alpha=28000)
    add_blob(slide, Inches(0.6), Inches(10.5), Inches(4.8), PINK, alpha=24000)
    add_blob(slide, Inches(7.0), Inches(12.0), Inches(3.8), ORANGE, alpha=20000)


def pill(slide, x, y, w, h, label, gradient_stops=None, fill=None, text_color=WHITE, size=11, letter_sp=400):
    if gradient_stops:
        add_gradient_rect(slide, x, y, w, h, gradient_stops, angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    else:
        add_rect(slide, x, y, w, h, fill or CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y, w, h, label, size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=letter_sp)


def wordmark(slide, x, y, size=18):
    tb = slide.shapes.add_textbox(x, y, Inches(2.2), Inches(0.45))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = 'Frame'
    r1.font.name = 'Helvetica Neue'; r1.font.size = Pt(size); r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = 'drops'
    r2.font.name = 'Helvetica Neue'; r2.font.size = Pt(size); r2.font.bold = True
    r2.font.color.rgb = WHITE
    add_gradient_rect(slide, x + Inches(1.85), y + Inches(0.13), Inches(0.18), Inches(0.18),
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=2700000, shape=MSO_SHAPE.OVAL)


def slide_chrome(slide, idx, total, duration):
    wordmark(slide, Inches(0.45), Inches(0.5), size=16)
    pw, ph = Inches(1.3), Inches(0.36)
    px = SW - pw - Inches(0.45)
    py = Inches(0.54)
    add_rect(slide, px, py, pw, ph, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, px, py, pw, ph, f'{idx:02d}/{total:02d} · {duration}s',
             size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=200)


def gradient_divider(slide, y, w_in=2.6, x_center=None):
    if x_center is None:
        x_center = SW // 2
    w = Inches(w_in); x = x_center - w // 2
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.06))
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', '0'); lin.set('scaled', '0')
    s.line.fill.background(); s.shadow.inherit = False


# Phone mockup
def phone_mockup(slide, x, y, w, h, screen_content=None):
    add_rect(slide, x, y, w, h, NAVY_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(slide, x, y, w, h,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=2.0, angle=2700000)
    inset = Inches(0.12)
    sx, sy = x + inset, y + inset
    sw, sh = w - 2 * inset, h - 2 * inset
    add_gradient_rect(slide, sx, sy, sw, sh,
                      stops=[(0.0, RGBColor(0x10, 0x1F, 0x42)),
                             (1.0, RGBColor(0x1A, 0x12, 0x44))],
                      angle=5400000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    nw = w * 0.32
    add_rect(slide, x + (w - nw) / 2, y + Inches(0.08), nw, Inches(0.18),
             NAVY_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if screen_content:
        screen_content(slide, sx, sy, sw, sh)


# Album-grid-with-hearts screen
def screen_gallery_hearts(slide, sx, sy, sw, sh):
    add_text(slide, sx + Inches(0.22), sy + Inches(0.38), sw - Inches(0.44), Inches(0.3),
             'Sneha & Rohan',
             size=11, bold=True, color=WHITE)
    add_text(slide, sx + Inches(0.22), sy + Inches(0.65), sw - Inches(0.44), Inches(0.25),
             '482 photos · Tap ♥ to pick',
             size=8, color=MUTED)
    cols = 3; rows = 5
    gap = Inches(0.07); pad = Inches(0.22)
    grid_top = sy + Inches(1.05)
    grid_w = sw - 2 * pad
    cw = (grid_w - (cols - 1) * gap) / cols
    ch = cw
    pal = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
        [(0.0, ROYAL), (1.0, PINK)],
        [(0.0, ORANGE), (1.0, PURPLE)],
        [(0.0, ROYAL), (1.0, ORANGE)],
    ]
    hearted = {(0, 1), (1, 0), (1, 2), (2, 1), (3, 2), (4, 0), (4, 1)}
    for r in range(rows):
        for c in range(cols):
            cx = sx + pad + c * (cw + gap)
            cy = grid_top + r * (ch + gap)
            add_gradient_rect(slide, cx, cy, cw, ch, pal[(r * cols + c) % len(pal)],
                              angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            if (r, c) in hearted:
                bd = Inches(0.28)
                badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                               cx + cw - bd - Inches(0.05),
                                               cy + ch - bd - Inches(0.05),
                                               bd, bd)
                badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
                badge.line.fill.background(); badge.shadow.inherit = False
                hh = slide.shapes.add_shape(MSO_SHAPE.HEART,
                                            cx + cw - bd - Inches(0.01),
                                            cy + ch - bd - Inches(0.01),
                                            bd * 0.78, bd * 0.7)
                hh.fill.solid(); hh.fill.fore_color.rgb = PINK
                hh.line.fill.background(); hh.shadow.inherit = False


# ============================================================
# 5 SLIDES — each labeled with target seconds
# ============================================================

TOTAL = 5
TIMINGS = [5, 6, 7, 6, 6]  # = 30s


# ---------- 1. HOOK (5s) ----------
def s1_hook():
    s = prs.slides.add_slide(BLANK)
    background(s)
    slide_chrome(s, 1, TOTAL, TIMINGS[0])

    pill(s, Inches(2.45), Inches(2.6), Inches(2.6), Inches(0.5),
         'FOR PHOTOGRAPHERS',
         gradient_stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)])

    add_text(s, Inches(0.55), Inches(3.7), Inches(6.4), Inches(1.2),
             'Still delivering',
             size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, Inches(0.4), Inches(4.9), Inches(6.7), Inches(3.6),
                      'photos\nthe hard\nway?',
                      size=96, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.0,
                      letter_sp=-30)

    # WhatsApp-chaos visual (chat bubbles spilling out)
    bubbles_y = Inches(9.6)
    bubbles = [
        ('"Send pic 87 again 🙏"',  Inches(0.6),  bubbles_y,          Inches(4.6), [(0.0, PINK), (1.0, ORANGE)]),
        ('"Wait, which one?"',      Inches(2.4),  bubbles_y + Inches(0.85), Inches(4.6), [(0.0, PURPLE), (1.0, PINK)]),
        ('"Where\'s the link?"',    Inches(0.4),  bubbles_y + Inches(1.7),  Inches(4.4), [(0.0, ROYAL), (1.0, PURPLE)]),
        ('"Resend 23 & 41 plz"',    Inches(2.6),  bubbles_y + Inches(2.55), Inches(4.4), [(0.0, ORANGE), (1.0, PINK)]),
    ]
    for txt, x, y, w, grad in bubbles:
        h = Inches(0.7)
        add_gradient_rect(s, x, y, w, h, grad, angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x + Inches(0.3), y, w - Inches(0.6), h,
                 txt, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ---------- 2. INTRO + WHAT IT DOES (6s) ----------
def s2_meet():
    s = prs.slides.add_slide(BLANK)
    background(s)
    slide_chrome(s, 2, TOTAL, TIMINGS[1])

    pill(s, Inches(2.85), Inches(2.0), Inches(1.8), Inches(0.5),
         'INTRODUCING', gradient_stops=[(0.0, PURPLE), (1.0, PINK)])

    add_text(s, 0, Inches(2.85), SW, Inches(0.8),
             'Meet',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(3.6), SW, Inches(2.0),
                      'Framedrops.',
                      size=98, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, letter_sp=-40)

    gradient_divider(s, Inches(5.85), w_in=2.4)

    add_text(s, Inches(0.6), Inches(6.2), Inches(6.3), Inches(1.0),
             'Upload. Share. Get picks.\nDeliver. Done.',
             size=20, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, line_spacing=1.4)

    # 4-step chip row with icons
    steps = [
        ('Upload',   [(0.0, ROYAL),  (1.0, PURPLE)],  '↑'),
        ('Share',    [(0.0, PURPLE), (1.0, PINK)],    '→'),
        ('Pick ♥',   [(0.0, PINK),   (1.0, ORANGE)],  '♥'),
        ('Deliver',  [(0.0, ORANGE), (1.0, ROYAL)],   '✓'),
    ]
    cy = Inches(8.0)
    cw = Inches(1.5); ch = Inches(1.7); gap = Inches(0.12)
    total_w = 4 * cw + 3 * gap
    sx = (SW - total_w) // 2
    for i, (label, grad, icon) in enumerate(steps):
        cx = sx + i * (cw + gap)
        add_rect(s, cx, cy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, cx, cy, cw, ch, grad, line_w_pt=1.5, angle=2700000)
        # icon circle
        ix = cx + (cw - Inches(0.7)) / 2
        iy = cy + Inches(0.25)
        add_gradient_rect(s, ix, iy, Inches(0.7), Inches(0.7), grad,
                          angle=2700000, shape=MSO_SHAPE.OVAL)
        add_text(s, ix, iy, Inches(0.7), Inches(0.7), icon,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx, cy + Inches(1.1), cw, Inches(0.45),
                 label, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=200)

    # mini benefit row
    by = Inches(10.2)
    benefits = ['Smart compression', 'One clean link', 'Works on any device']
    bw = Inches(6.4); bh = Inches(2.4)
    bx = (SW - bw) // 2
    add_rect(s, bx, by, bw, bh, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, bx, by, bw, bh,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=1.25, angle=2700000)
    for i, b in enumerate(benefits):
        ry = by + Inches(0.25) + i * Inches(0.65)
        add_gradient_rect(s, bx + Inches(0.35), ry + Inches(0.18), Inches(0.22), Inches(0.22),
                          stops=[(0.0, ROYAL), (1.0, PINK)], angle=0, shape=MSO_SHAPE.OVAL)
        add_text(s, bx + Inches(0.35), ry + Inches(0.18), Inches(0.22), Inches(0.22),
                 '✓', size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx + Inches(0.75), ry, bw - Inches(1.0), Inches(0.5),
                 b, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ---------- 3. THE MAGIC — Clients heart, you transfer (7s) ----------
def s3_magic():
    s = prs.slides.add_slide(BLANK)
    background(s)
    slide_chrome(s, 3, TOTAL, TIMINGS[2])

    pill(s, Inches(2.45), Inches(1.7), Inches(2.6), Inches(0.5),
         'THE MAGIC ✦', gradient_stops=[(0.0, PINK), (1.0, ORANGE)])

    add_text(s, 0, Inches(2.55), SW, Inches(0.8),
             'Clients',
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(3.3), SW, Inches(2.0),
                      'heart it.',
                      size=82, bold=True,
                      stops=[(0.0, PINK), (0.5, PURPLE), (1.0, ROYAL)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    # phone with gallery + hearts
    phone_w = Inches(3.0); phone_h = Inches(5.0)
    phone_x = (SW - phone_w) // 2
    phone_y = Inches(5.5)
    phone_mockup(s, phone_x, phone_y, phone_w, phone_h, screen_content=screen_gallery_hearts)

    # arrow down
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                             SW / 2 - Inches(0.5), Inches(10.75),
                             Inches(1.0), Inches(0.55))
    sp = arr.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', '5400000'); lin.set('scaled', '0')
    arr.line.fill.background(); arr.shadow.inherit = False

    # folder card — local transfer
    cy = Inches(11.5)
    cw = Inches(6.4); ch = Inches(1.3)
    cx = (SW - cw) // 2
    add_rect(s, cx, cy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, cx, cy, cw, ch,
                           stops=[(0.0, ORANGE), (1.0, PINK)], line_w_pt=1.5, angle=2700000)
    # folder icon
    fx = cx + Inches(0.3); fy = cy + Inches(0.35)
    add_gradient_rect(s, fx, fy, Inches(0.3), Inches(0.12),
                      stops=[(0.0, ORANGE), (1.0, PINK)], angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_rect(s, fx - Inches(0.02), fy + Inches(0.08), Inches(0.75), Inches(0.55),
                      stops=[(0.0, ORANGE), (1.0, PINK)], angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, cx + Inches(1.2), cy + Inches(0.2), cw - Inches(1.4), Inches(0.4),
             'Originals transfer locally',
             size=14, bold=True, color=WHITE)
    add_text(s, cx + Inches(1.2), cy + Inches(0.65), cw - Inches(1.4), Inches(0.4),
             '/Sneha-Rohan/Selected  · 124 files',
             size=10, color=MUTED)


# ---------- 4. OUTCOME (6s) ----------
def s4_outcome():
    s = prs.slides.add_slide(BLANK)
    background(s)
    slide_chrome(s, 4, TOTAL, TIMINGS[3])

    pill(s, Inches(2.85), Inches(2.0), Inches(1.8), Inches(0.5),
         'THE OUTCOME', gradient_stops=[(0.0, ORANGE), (1.0, PINK)])

    add_text(s, 0, Inches(3.0), SW, Inches(0.9),
             'Less',
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(3.8), SW, Inches(1.8),
                      'chaos.',
                      size=96, bold=True,
                      stops=[(0.0, PINK), (0.5, ORANGE), (1.0, PURPLE)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    add_text(s, 0, Inches(5.85), SW, Inches(0.9),
             'More',
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(6.65), SW, Inches(1.8),
                      'shoots.',
                      size=96, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    add_text(s, 0, Inches(8.7), SW, Inches(0.9),
             'More',
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(9.5), SW, Inches(1.8),
                      'business.',
                      size=96, bold=True,
                      stops=[(0.0, PURPLE), (0.5, PINK), (1.0, ORANGE)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    gradient_divider(s, Inches(11.7), w_in=2.4)
    add_text(s, 0, Inches(12.0), SW, Inches(0.4),
             'Built for photographers who move fast.',
             size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=300)


# ---------- 5. CTA (6s) ----------
def s5_cta():
    s = prs.slides.add_slide(BLANK)
    background(s)
    slide_chrome(s, 5, TOTAL, TIMINGS[4])

    pill(s, Inches(2.95), Inches(2.8), Inches(1.6), Inches(0.5),
         'START NOW', gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)])

    add_text(s, 0, Inches(3.8), SW, Inches(0.9),
             'Try it today',
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_gradient_text(s, 0, Inches(5.0), SW, Inches(3.0),
                      'framedrops\n.in',
                      size=104, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.05,
                      letter_sp=-50)

    # CTA button
    bw, bh = Inches(5.0), Inches(1.1)
    bx = (SW - bw) // 2
    by = Inches(8.7)
    add_gradient_rect(s, bx, by, bw, bh,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, bx, by, bw, bh, 'Start free  →',
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=200)

    # 3 USP chips
    usps = [
        ('Free to start',  [(0.0, ROYAL),  (1.0, PURPLE)]),
        ('No card',        [(0.0, PURPLE), (1.0, PINK)]),
        ('2 min setup',    [(0.0, PINK),   (1.0, ORANGE)]),
    ]
    uy = Inches(10.2)
    cw = Inches(2.0); ch = Inches(1.1); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    sx = (SW - total_w) // 2
    for i, (label, grad) in enumerate(usps):
        cx = sx + i * (cw + gap)
        add_rect(s, cx, uy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, cx, uy, cw, ch, grad, line_w_pt=1.25, angle=2700000)
        add_gradient_rect(s, cx + (cw - Inches(0.2)) / 2, uy + Inches(0.18), Inches(0.2), Inches(0.2),
                          stops=grad, angle=2700000, shape=MSO_SHAPE.OVAL)
        add_text(s, cx, uy + Inches(0.5), cw, Inches(0.55),
                 label, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 0, Inches(11.8), SW, Inches(0.4),
             '@framedrops · framedrops.in',
             size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=400)
    gradient_divider(s, Inches(12.4), w_in=2.0)


# Build
s1_hook()
s2_meet()
s3_magic()
s4_outcome()
s5_cta()

out_path = '/Users/apple/Desktop/PHOTOSHARE/framedropnotes/Framedrops_Reel_30s.pptx'
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Aspect: {prs.slide_width / prs.slide_height:.4f} (target 9/16 = {9/16:.4f})')
print(f'Slides: {len(prs.slides)}  ·  Total: {sum(TIMINGS)}s')
print(f'Per-slide timing: {TIMINGS}')
