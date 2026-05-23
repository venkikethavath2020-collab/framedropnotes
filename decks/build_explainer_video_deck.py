"""
Framedrops 30-Second Explainer Video Deck — 4 slides. MODERN edition.

Vertical 9:16 (1080×1920 equivalent) — phone-shaped. WhatsApp Status /
Instagram Reels / YouTube Shorts ready.

Design philosophy this time:
  - Mostly white / off-white backgrounds (Linear, Notion, Stripe vibe)
  - Single brand accent: purple #7C3AED
  - Generous whitespace
  - Black/dark text (high readability on small phone screens)
  - NO heavy gradients. One subtle pill accent maximum.
  - Modern sans-serif feel (Inter), product-grade typography

Story:
  1. THE CHAOS         — Drive + paper + WhatsApp (red, restrained)
  2. CREATE & UPLOAD   — Drop photos into Framedrops
  3. SHARE & SELECT    — Client taps hearts on favorites
  4. LOCAL TRANSFER    — Copy originals locally + CTA URL

Run:
    python3 build_explainer_video_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ─── Restrained palette — Linear / Notion / Stripe inspired ────────────────
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)   # single brand accent
PURPLE_DARK  = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_TINT  = RGBColor(0xF5, 0xF3, 0xFF)   # subtle background tint
PURPLE_LINE  = RGBColor(0xDD, 0xD6, 0xFE)

INK          = RGBColor(0x0F, 0x17, 0x2A)   # near-black text
SLATE_700    = RGBColor(0x33, 0x41, 0x55)   # body text
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)   # secondary text
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)   # muted
SLATE_300    = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)   # subtle borders
SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)   # very subtle bg
SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)   # off-white base

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xFA, 0xFA, 0xFA)

# Functional accents — used sparingly
RED          = RGBColor(0xDC, 0x26, 0x26)   # chaos slide only
RED_TINT     = RGBColor(0xFE, 0xF2, 0xF2)
GREEN        = RGBColor(0x05, 0x96, 0x69)   # success state
GREEN_TINT   = RGBColor(0xEC, 0xFD, 0xF5)

# 9:16 portrait — vertical phone-shaped
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.33)

OUT_FILE = Path(__file__).parent / "FrameDrops_Explainer_Video.pptx"
WEBSITE  = "https://framedrops.in/"


# ─── Helpers ──────────────────────────────────────────────────────────────


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rgb_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def apply_gradient_fill(shape, stops, angle: float = 90) -> None:
    """Used VERY sparingly — only for the single brand pill per slide."""
    sp_pr = shape._element.spPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        for el in sp_pr.findall(qn(tag)):
            sp_pr.remove(el)
    grad = etree.SubElement(sp_pr, qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(int(pos * 1000)))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", _rgb_hex(color))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "0")


def solid_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def add_solid_card(slide, left, top, width, height, fill,
                   corner=0.04, border=None, border_width=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card


def add_text(
    slide, left, top, width, height, text,
    *, size=24, color=INK, bold=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    line_spacing=1.2, font="Inter",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_brand_pill(slide, left, top, width, height, text,
                    fill=PURPLE, fg=WHITE, size=11):
    """The ONE accent element — small, subtle pill."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    tf = pill.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Inter"
    return pill


def add_footer(slide, on_dark=False):
    """Clean URL footer — no chrome, just the URL."""
    color = SLATE_500 if not on_dark else SLATE_300
    # Thin separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(12.5),
                                   Inches(6.5), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = SLATE_200 if not on_dark else SLATE_700
    sep.line.fill.background()
    # URL — small, clean
    add_text(slide, Inches(0.5), Inches(12.7), Inches(6.5), Inches(0.4),
             WEBSITE,
             size=14, color=color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_step_indicator(slide, current, total=4):
    """Top-right step indicator — minimal."""
    add_text(slide, Inches(5.5), Inches(0.6), Inches(1.8), Inches(0.3),
             f"{current} / {total}",
             size=11, color=SLATE_400, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_brand_mark(slide):
    """Top-left brand mark — minimal."""
    # Small purple square
    sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.5), Inches(0.55),
                                  Inches(0.32), Inches(0.32))
    sq.adjustments[0] = 0.2
    sq.fill.solid()
    sq.fill.fore_color.rgb = PURPLE
    sq.line.fill.background()
    # Wordmark
    add_text(slide, Inches(0.95), Inches(0.55), Inches(3), Inches(0.35),
             "Framedrops",
             size=13, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_thumbnail_grid(slide, x_in, y_in, cols, rows, cell_w_in, cell_h_in,
                        gap_in=0.06, hearts=None, hint_colors=None):
    """A clean photo-grid mock. Cells are neutral grey by default with
    optional varied tint to suggest "photos" without being garish.
    `hearts` is a set of (row, col) to mark as favorites."""
    hearts = hearts or set()
    # Subtle tints — barely distinguishable, gives "photo grid" feel
    tints = hint_colors or [
        RGBColor(0xE5, 0xE7, 0xEB),
        RGBColor(0xDC, 0xDE, 0xE3),
        RGBColor(0xE9, 0xE5, 0xF1),
        RGBColor(0xE4, 0xE7, 0xEB),
        RGBColor(0xEB, 0xE9, 0xEE),
        RGBColor(0xE0, 0xE3, 0xE7),
    ]
    for r in range(rows):
        for c in range(cols):
            tx = Inches(x_in + c * (cell_w_in + gap_in))
            ty = Inches(y_in + r * (cell_h_in + gap_in))
            tint = tints[(r * cols + c) % len(tints)]
            add_solid_card(slide, tx, ty, Inches(cell_w_in), Inches(cell_h_in),
                           tint, corner=0.08)
            # Tiny "image" placeholder hint (centered dot) — keeps it modern
            dot_size = 0.18
            add_solid_card(slide,
                           tx + Inches(cell_w_in/2 - dot_size/2),
                           ty + Inches(cell_h_in/2 - dot_size/2),
                           Inches(dot_size), Inches(dot_size),
                           SLATE_300, corner=0.5)
            # Heart marker for favorited cells
            if (r, c) in hearts:
                heart_size = 0.32
                hx = tx + Inches(cell_w_in - heart_size - 0.08)
                hy = ty + Inches(0.08)
                # White circle bg
                add_solid_card(slide, hx, hy,
                               Inches(heart_size), Inches(heart_size),
                               WHITE, corner=0.5)
                # Filled heart shape
                heart_pad = 0.05
                h = slide.shapes.add_shape(MSO_SHAPE.HEART,
                                            hx + Inches(heart_pad),
                                            hy + Inches(heart_pad + 0.01),
                                            Inches(heart_size - 2*heart_pad),
                                            Inches(heart_size - 2*heart_pad))
                h.fill.solid()
                h.fill.fore_color.rgb = PURPLE
                h.line.fill.background()


def add_phone_mock(slide, x_in, y_in, w_in, h_in):
    """Draw a phone-frame outer shell. Returns inner content rect bounds."""
    # Outer shell (thin dark border)
    bezel = 0.12
    add_solid_card(slide, Inches(x_in - bezel), Inches(y_in - bezel),
                   Inches(w_in + 2*bezel), Inches(h_in + 2*bezel),
                   INK, corner=0.08)
    # Inner screen
    add_solid_card(slide, Inches(x_in), Inches(y_in),
                   Inches(w_in), Inches(h_in),
                   WHITE, corner=0.05)
    return x_in, y_in, w_in, h_in


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — THE CHAOS
# ────────────────────────────────────────────────────────────────────────────


def slide_1_chaos(prs):
    s = blank(prs)
    solid_background(s, SLATE_50)

    add_brand_mark(s)
    add_step_indicator(s, 1)

    # Eyebrow
    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(0.4),
             "THE PROBLEM",
             size=11, color=RED, bold=True,
             align=PP_ALIGN.LEFT)

    # Headline
    add_text(s, Inches(0.5), Inches(2.1), Inches(6.5), Inches(2.4),
             "Photo selection\nis broken.",
             size=52, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # Subhead
    add_text(s, Inches(0.5), Inches(4.6), Inches(6.5), Inches(0.6),
             "Hours wasted on every wedding.",
             size=18, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    # Three pain cards — clean, minimal, white with subtle border
    pains = [
        ("📁", "Google Drive", "30-min uploads."),
        ("📝", "Paper chits", "Decode handwriting."),
        ("💬", "WhatsApp mess", "200+ confusing messages."),
    ]
    y = 6.0
    for emoji, title, body in pains:
        # White card with subtle border
        add_solid_card(s, Inches(0.5), Inches(y), Inches(6.5), Inches(1.7),
                       WHITE, corner=0.05,
                       border=SLATE_200, border_width=1)
        # Emoji
        add_text(s, Inches(0.85), Inches(y + 0.5), Inches(1.0), Inches(1.0),
                 emoji, size=36, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, Inches(2.0), Inches(y + 0.4), Inches(4.5), Inches(0.5),
                 title, size=20, color=INK, bold=True,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        # Body
        add_text(s, Inches(2.0), Inches(y + 0.95), Inches(4.5), Inches(0.5),
                 body, size=14, color=SLATE_500,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        y += 1.95

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — CREATE ALBUM + UPLOAD
# ────────────────────────────────────────────────────────────────────────────


def slide_2_upload(prs):
    s = blank(prs)
    solid_background(s, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 2)

    # Eyebrow
    add_brand_pill(s, Inches(0.5), Inches(1.6), Inches(1.0), Inches(0.36),
                    "STEP 1", fill=PURPLE, fg=WHITE, size=10)

    # Headline
    add_text(s, Inches(0.5), Inches(2.2), Inches(6.5), Inches(2.4),
             "Upload\nyour photos.",
             size=52, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # Subhead
    add_text(s, Inches(0.5), Inches(4.7), Inches(6.5), Inches(0.6),
             "Drag, drop, done. Compressed in browser.",
             size=18, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    # Phone mock — centered
    phone_x = 1.6
    phone_y = 5.8
    phone_w = 4.3
    phone_h = 5.7
    add_phone_mock(s, phone_x, phone_y, phone_w, phone_h)

    # In-phone content
    pad = 0.3

    # Top: album name
    add_text(s, Inches(phone_x + pad), Inches(phone_y + 0.3),
             Inches(phone_w - 2*pad), Inches(0.3),
             "STUDIO · RAVI + ANJALI",
             size=8, color=PURPLE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(phone_x + pad), Inches(phone_y + 0.6),
             Inches(phone_w - 2*pad), Inches(0.5),
             "Wedding Day 1",
             size=18, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Progress bar
    progress_y = phone_y + 1.4
    # Track
    add_solid_card(s, Inches(phone_x + pad), Inches(progress_y),
                   Inches(phone_w - 2*pad), Inches(0.16),
                   SLATE_100, corner=0.5)
    # Fill (78%)
    add_solid_card(s, Inches(phone_x + pad), Inches(progress_y),
                   Inches((phone_w - 2*pad) * 0.78), Inches(0.16),
                   PURPLE, corner=0.5)
    # Label below
    add_text(s, Inches(phone_x + pad), Inches(progress_y + 0.25),
             Inches(phone_w - 2*pad), Inches(0.3),
             "Uploading · 78%",
             size=10, color=SLATE_500, bold=True,
             align=PP_ALIGN.LEFT)

    # Grid 3×3
    grid_y = phone_y + 2.3
    cell_w = (phone_w - 2*pad - 2*0.06) / 3
    add_thumbnail_grid(s,
                       x_in=phone_x + pad,
                       y_in=grid_y,
                       cols=3, rows=3,
                       cell_w_in=cell_w,
                       cell_h_in=0.92,
                       gap_in=0.06)

    # Footer caption
    add_text(s, Inches(0.5), Inches(11.7), Inches(6.5), Inches(0.5),
             "25 MB → 250 KB. Auto-compressed.",
             size=14, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — SHARE LINK + CLIENT SELECTS
# ────────────────────────────────────────────────────────────────────────────


def slide_3_selects(prs):
    s = blank(prs)
    solid_background(s, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 3)

    add_brand_pill(s, Inches(0.5), Inches(1.6), Inches(1.0), Inches(0.36),
                    "STEP 2", fill=PURPLE, fg=WHITE, size=10)

    add_text(s, Inches(0.5), Inches(2.2), Inches(6.5), Inches(2.4),
             "Client picks\nfavorites.",
             size=52, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    add_text(s, Inches(0.5), Inches(4.7), Inches(6.5), Inches(0.6),
             "One link via WhatsApp. They tap a heart.",
             size=18, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    # Phone mock
    phone_x = 1.6
    phone_y = 5.8
    phone_w = 4.3
    phone_h = 5.7
    add_phone_mock(s, phone_x, phone_y, phone_w, phone_h)

    pad = 0.3

    # Top
    add_text(s, Inches(phone_x + pad), Inches(phone_y + 0.3),
             Inches(phone_w - 2*pad), Inches(0.3),
             "RAVI + ANJALI · GALLERY",
             size=8, color=PURPLE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(phone_x + pad), Inches(phone_y + 0.6),
             Inches(phone_w - 2*pad), Inches(0.5),
             "Pick your favorites",
             size=18, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    # Counter pill
    add_solid_card(s, Inches(phone_x + pad), Inches(phone_y + 1.2),
                   Inches(2.0), Inches(0.34),
                   PURPLE_TINT, corner=0.5,
                   border=PURPLE_LINE, border_width=1)
    add_text(s, Inches(phone_x + pad), Inches(phone_y + 1.2),
             Inches(2.0), Inches(0.34),
             "❤ 5 of 1,500 picked",
             size=10, color=PURPLE_DARK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Grid 3×4 with hearts
    grid_y = phone_y + 1.8
    cell_w = (phone_w - 2*pad - 2*0.06) / 3
    hearts = {(0, 1), (1, 0), (1, 2), (2, 1), (3, 2)}
    add_thumbnail_grid(s,
                       x_in=phone_x + pad,
                       y_in=grid_y,
                       cols=3, rows=4,
                       cell_w_in=cell_w,
                       cell_h_in=0.88,
                       gap_in=0.06,
                       hearts=hearts)

    # Footer caption
    add_text(s, Inches(0.5), Inches(11.7), Inches(6.5), Inches(0.5),
             "No paper. No decoding. Real-time count.",
             size=14, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — LOCAL TRANSFER + CTA
# ────────────────────────────────────────────────────────────────────────────


def slide_4_transfer(prs):
    s = blank(prs)
    solid_background(s, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 4)

    add_brand_pill(s, Inches(0.5), Inches(1.6), Inches(1.0), Inches(0.36),
                    "STEP 3", fill=PURPLE, fg=WHITE, size=10)

    add_text(s, Inches(0.5), Inches(2.2), Inches(6.5), Inches(2.4),
             "Copy originals\nlocally.",
             size=52, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    add_text(s, Inches(0.5), Inches(4.7), Inches(6.5), Inches(0.6),
             "Selected list → your laptop. No upload needed.",
             size=18, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    # Two folders + arrow — clean, minimal
    folder_y = 6.0
    folder_w = 2.4
    folder_h = 2.0

    # Source folder card
    add_solid_card(s, Inches(0.5), Inches(folder_y),
                   Inches(folder_w), Inches(folder_h),
                   SLATE_50, corner=0.06,
                   border=SLATE_200, border_width=1)
    add_text(s, Inches(0.5), Inches(folder_y + 0.4),
             Inches(folder_w), Inches(0.7),
             "📁", size=44, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(folder_y + 1.2),
             Inches(folder_w), Inches(0.35),
             "All photos",
             size=14, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.5), Inches(folder_y + 1.55),
             Inches(folder_w), Inches(0.3),
             "Your laptop",
             size=11, color=SLATE_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Arrow — minimal
    arrow_x = 3.05
    arrow_w = 1.4
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(arrow_x), Inches(folder_y + 0.85),
                                Inches(arrow_w), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = PURPLE
    arrow.line.fill.background()
    # "200 selected" label above
    add_text(s, Inches(arrow_x - 0.1), Inches(folder_y + 0.35),
             Inches(arrow_w + 0.2), Inches(0.4),
             "200 selected",
             size=11, color=PURPLE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # "Local copy" label below
    add_text(s, Inches(arrow_x - 0.1), Inches(folder_y + 1.4),
             Inches(arrow_w + 0.2), Inches(0.4),
             "Local copy",
             size=11, color=SLATE_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Destination folder card — subtle green tint (success state)
    dest_x = 4.6
    add_solid_card(s, Inches(dest_x), Inches(folder_y),
                   Inches(folder_w), Inches(folder_h),
                   GREEN_TINT, corner=0.06,
                   border=GREEN, border_width=1)
    add_text(s, Inches(dest_x), Inches(folder_y + 0.4),
             Inches(folder_w), Inches(0.7),
             "✓", size=44, color=GREEN, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(dest_x), Inches(folder_y + 1.2),
             Inches(folder_w), Inches(0.35),
             "Favorites",
             size=14, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(dest_x), Inches(folder_y + 1.55),
             Inches(folder_w), Inches(0.3),
             "Ready to deliver",
             size=11, color=SLATE_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Quiet note about no internet upload
    add_text(s, Inches(0.5), Inches(8.4), Inches(6.5), Inches(0.5),
             "Browser file copy. Bytes never leave your laptop.",
             size=13, color=SLATE_500,
             align=PP_ALIGN.CENTER)

    # CTA card — clean, purple accent, prominent
    cta_y = 9.5
    add_solid_card(s, Inches(0.5), Inches(cta_y),
                   Inches(6.5), Inches(2.4),
                   PURPLE, corner=0.05)
    # Free badge — sits inside, small
    add_solid_card(s, Inches(0.5 + 0.4), Inches(cta_y + 0.4),
                   Inches(1.6), Inches(0.36),
                   WHITE, corner=0.5)
    add_text(s, Inches(0.5 + 0.4), Inches(cta_y + 0.4),
             Inches(1.6), Inches(0.36),
             "300 FREE",
             size=10, color=PURPLE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5 + 0.4), Inches(cta_y + 0.9),
             Inches(6.5 - 0.8), Inches(0.8),
             "Try Framedrops free",
             size=28, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    add_text(s, Inches(0.5 + 0.4), Inches(cta_y + 1.7),
             Inches(6.5 - 0.8), Inches(0.4),
             "No credit card · 2-minute signup",
             size=13, color=PURPLE_TINT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────


def build() -> None:
    prs = make_presentation()
    slide_1_chaos(prs)
    slide_2_upload(prs)
    slide_3_selects(prs)
    slide_4_transfer(prs)
    prs.save(OUT_FILE)
    print(f"✓ Wrote {OUT_FILE}")
    print(f"  4 slides @ 9:16 vertical, modern design (white + purple accent)")


if __name__ == "__main__":
    build()
