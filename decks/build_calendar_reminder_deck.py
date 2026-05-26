"""
Framedrops Calendar Reminder — 3-slide portrait WhatsApp share.

Vertical 9:16 (1080×1920 equivalent). One slide per beat:
  1. THE PAIN  — "Forgot another shoot date?" + scribbled calendar mock
  2. THE FIX   — "Add it to your Framedrops calendar." + calendar +
                  email-notification mock
  3. THE CTA   — "Try Framedrops free." + phone-screen mock + big link pill

Design language:
  - Off-white base with subtle tinted backdrops per slide
  - Multi-hue accent palette (purple primary, amber + emerald + indigo
    used sparingly to give each slide its own energy without going garish)
  - Slate text, generous whitespace, soft shadows on cards
  - All mocks drawn in code — no external images

Run:
    python3 build_calendar_reminder_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ─── Palette — modern, multi-hue but restrained ─────────────────────────────
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)   # brand primary
PURPLE_DARK  = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_TINT  = RGBColor(0xF5, 0xF3, 0xFF)
PURPLE_LINE  = RGBColor(0xDD, 0xD6, 0xFE)

INDIGO       = RGBColor(0x4F, 0x46, 0xE5)   # calendar accents
INDIGO_TINT  = RGBColor(0xEE, 0xF2, 0xFF)

AMBER        = RGBColor(0xF5, 0x9E, 0x0B)   # warning / pain slide
AMBER_TINT   = RGBColor(0xFF, 0xFB, 0xEB)

EMERALD      = RGBColor(0x05, 0x96, 0x69)   # success / notification dot
EMERALD_TINT = RGBColor(0xEC, 0xFD, 0xF5)

PINK         = RGBColor(0xEC, 0x48, 0x99)   # subtle 4th accent for calendar event chips
PINK_TINT    = RGBColor(0xFD, 0xF2, 0xF8)

INK          = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700    = RGBColor(0x33, 0x41, 0x55)
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300    = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xDC, 0x26, 0x26)   # "missed" red dot only

# 9:16 portrait — WhatsApp-ready
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.33)

OUT_FILE = Path(__file__).parent / "FrameDrops_Calendar_Reminder.pptx"
WEBSITE  = "https://framedrops.in/"

FONT = "Inter"


# ─── Helpers ────────────────────────────────────────────────────────────────


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rgb_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def apply_gradient_fill(shape, stops, angle: float = 90) -> None:
    """Soft two-stop gradient used sparingly for slide backdrops."""
    sp_pr = shape._element.spPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        for el in sp_pr.findall(qn(tag)):
            sp_pr.remove(el)
    grad = etree.SubElement(sp_pr, qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(int(pos * 1000)))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", _rgb_hex(color))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "0")


def solid_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def gradient_background(slide, top_color, bottom_color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    apply_gradient_fill(bg, [(0.0, top_color), (1.0, bottom_color)], angle=90)


def add_card(slide, left, top, width, height, fill,
             corner=0.04, border=None, border_width=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card


def add_text(
    slide, left, top, width, height, text,
    *, size=24, color=INK, bold=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    line_spacing=1.2, font=FONT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_brand_mark(slide):
    sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.5), Inches(0.55),
                                  Inches(0.32), Inches(0.32))
    sq.adjustments[0] = 0.2
    sq.fill.solid()
    sq.fill.fore_color.rgb = PURPLE
    sq.line.fill.background()
    add_text(slide, Inches(0.95), Inches(0.55), Inches(3), Inches(0.35),
             "Framedrops",
             size=13, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_step_indicator(slide, current, total=3):
    add_text(slide, Inches(5.5), Inches(0.6), Inches(1.8), Inches(0.3),
             f"{current} / {total}",
             size=11, color=SLATE_400, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, color=SLATE_500):
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(12.5),
                                   Inches(6.5), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = SLATE_200
    sep.line.fill.background()
    add_text(slide, Inches(0.5), Inches(12.7), Inches(6.5), Inches(0.4),
             WEBSITE,
             size=14, color=color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_eyebrow(slide, top_in, text, color):
    add_text(slide, Inches(0.5), Inches(top_in), Inches(6.5), Inches(0.4),
             text,
             size=11, color=color, bold=True,
             align=PP_ALIGN.LEFT)


def add_dot(slide, x_in, y_in, size_in, color):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x_in), Inches(y_in),
                                Inches(size_in), Inches(size_in))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()
    return d


# ─── Mocks ──────────────────────────────────────────────────────────────────


def draw_calendar_mock(slide, x_in, y_in, w_in, h_in,
                       header_color=PURPLE,
                       events=None,
                       miss_day=None):
    """Mini monthly-calendar mock. 7 cols × 5 rows of date cells.
    `events`: dict { day_index_0_to_34: accent_color }
    `miss_day`: int — single day rendered with a red "missed" dot.
    """
    events = events or {}

    # Card shell
    add_card(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
             WHITE, corner=0.05, border=SLATE_200, border_width=0.75)

    # Soft shadow strip below the card (subtle)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x_in + 0.05),
                                     Inches(y_in + h_in - 0.04),
                                     Inches(w_in - 0.10),
                                     Inches(0.12))
    shadow.adjustments[0] = 0.5
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SLATE_200
    shadow.line.fill.background()

    # Header bar (colored strip with month name)
    header_h = 0.55
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(header_h))
    hdr.adjustments[0] = 0.05
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    add_text(slide, Inches(x_in + 0.3), Inches(y_in + 0.04),
             Inches(w_in - 0.6), Inches(header_h - 0.08),
             "MARCH 2026",
             size=11, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Dot navigators on the right side of the header
    nav_y = y_in + header_h / 2 - 0.06
    for offset, _ in enumerate([0, 1]):
        add_dot(slide, x_in + w_in - 0.4 - offset * 0.22, nav_y, 0.12,
                RGBColor(0xFF, 0xFF, 0xFF))

    # Weekday strip
    wd_y = y_in + header_h + 0.08
    weekdays = ["S", "M", "T", "W", "T", "F", "S"]
    cell_w = (w_in - 0.4) / 7
    for i, wd in enumerate(weekdays):
        add_text(slide, Inches(x_in + 0.2 + i * cell_w),
                 Inches(wd_y), Inches(cell_w), Inches(0.25),
                 wd, size=9, color=SLATE_400, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Date grid (5 rows × 7 cols)
    grid_y = wd_y + 0.32
    grid_h = h_in - (grid_y - y_in) - 0.25
    cell_h = grid_h / 5

    for row in range(5):
        for col in range(7):
            idx = row * 7 + col
            day_num = idx + 1  # naive — we just need visual fill, not real cal
            cx = x_in + 0.2 + col * cell_w
            cy = grid_y + row * cell_h

            # Event-day cells get a tinted square behind the number
            if idx in events:
                accent = events[idx]
                cell_pad = 0.05
                add_card(slide,
                         Inches(cx + cell_pad), Inches(cy + cell_pad),
                         Inches(cell_w - 2 * cell_pad),
                         Inches(cell_h - 2 * cell_pad),
                         accent, corner=0.15)
                # white text for filled cells
                add_text(slide, Inches(cx), Inches(cy),
                         Inches(cell_w), Inches(cell_h),
                         str(day_num),
                         size=10, color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                # Plain date number
                color = SLATE_300 if day_num > 28 else SLATE_700
                add_text(slide, Inches(cx), Inches(cy),
                         Inches(cell_w), Inches(cell_h),
                         str(day_num),
                         size=10, color=color,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # Red "missed" dot on a specific day
            if miss_day is not None and idx == miss_day:
                add_dot(slide, cx + cell_w - 0.18, cy + 0.04, 0.10, RED)


def draw_email_notification(slide, x_in, y_in, w_in, h_in):
    """Inbox notification mock — envelope icon + sender + subject + preview."""
    # Card
    add_card(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
             WHITE, corner=0.04, border=SLATE_200, border_width=0.75)

    # Sender avatar (purple square w/ F)
    av_size = 0.55
    av = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x_in + 0.25),
                                 Inches(y_in + 0.25),
                                 Inches(av_size), Inches(av_size))
    av.adjustments[0] = 0.2
    av.fill.solid()
    av.fill.fore_color.rgb = PURPLE
    av.line.fill.background()
    add_text(slide,
             Inches(x_in + 0.25), Inches(y_in + 0.25),
             Inches(av_size), Inches(av_size),
             "F", size=18, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Sender + time row
    add_text(slide, Inches(x_in + 0.95), Inches(y_in + 0.25),
             Inches(w_in - 1.5), Inches(0.3),
             "Framedrops Reminders",
             size=12, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(x_in + w_in - 0.95), Inches(y_in + 0.25),
             Inches(0.7), Inches(0.3),
             "9:00 AM",
             size=10, color=SLATE_400,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Subject
    add_text(slide, Inches(x_in + 0.95), Inches(y_in + 0.6),
             Inches(w_in - 1.1), Inches(0.4),
             "Tomorrow: Sharma Wedding shoot",
             size=13, color=INK, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Preview line
    add_text(slide, Inches(x_in + 0.95), Inches(y_in + 1.0),
             Inches(w_in - 1.1), Inches(0.35),
             "10 AM · Hyatt Hyderabad · 6 hrs",
             size=11, color=SLATE_500,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Unread dot
    add_dot(slide, x_in + 0.05, y_in + h_in / 2 - 0.06, 0.12, EMERALD)


def draw_phone_mock(slide, x_in, y_in, w_in, h_in):
    """Tilted phone-screen mock with the app's calendar visible. Used on CTA."""
    bezel = 0.10
    # Phone shell
    add_card(slide, Inches(x_in - bezel), Inches(y_in - bezel),
             Inches(w_in + 2 * bezel), Inches(h_in + 2 * bezel),
             INK, corner=0.08)
    # Screen
    add_card(slide, Inches(x_in), Inches(y_in),
             Inches(w_in), Inches(h_in),
             WHITE, corner=0.05)
    # Status bar dot
    add_dot(slide, x_in + w_in / 2 - 0.06, y_in + 0.12, 0.10, INK)

    # Mini calendar inside the phone
    cal_pad = 0.25
    draw_calendar_mock(slide,
                       x_in + cal_pad,
                       y_in + 0.55,
                       w_in - 2 * cal_pad,
                       h_in - 1.4,
                       header_color=PURPLE,
                       events={9: INDIGO, 14: EMERALD, 22: AMBER, 27: PINK})

    # Bottom CTA inside phone
    btn_w = w_in - 2 * cal_pad
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x_in + cal_pad),
                                  Inches(y_in + h_in - 0.75),
                                  Inches(btn_w), Inches(0.5))
    btn.adjustments[0] = 0.4
    btn.fill.solid()
    btn.fill.fore_color.rgb = PURPLE
    btn.line.fill.background()
    add_text(slide,
             Inches(x_in + cal_pad), Inches(y_in + h_in - 0.75),
             Inches(btn_w), Inches(0.5),
             "+ Add event",
             size=12, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_link_pill(slide, top_in, label):
    pill_w = 4.8
    pill_h = 1.05
    pill_x = (7.5 - pill_w) / 2
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(pill_x), Inches(top_in),
                                    Inches(pill_w), Inches(pill_h))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = PURPLE
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — THE PAIN
# ────────────────────────────────────────────────────────────────────────────


def slide_1_pain(prs):
    s = blank(prs)
    gradient_background(s, AMBER_TINT, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 1)

    add_eyebrow(s, 1.5, "THE PROBLEM", color=AMBER)

    # Headline
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.5), Inches(2.6),
             "Forgot another\nshoot date?",
             size=54, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # Subhead
    add_text(s, Inches(0.5), Inches(4.75), Inches(6.5), Inches(0.6),
             "Late to the venue. Awkward client call. Lost booking.",
             size=17, color=SLATE_500, line_spacing=1.4,
             align=PP_ALIGN.LEFT)

    # Visual — a calendar with one missed date highlighted in red
    draw_calendar_mock(s, 0.75, 5.9, 6.0, 5.7,
                       header_color=AMBER,
                       events={
                           4: SLATE_300,    # past/empty event placeholder
                           11: SLATE_300,
                           18: SLATE_300,
                       },
                       miss_day=14)  # the missed shoot

    # Caption under calendar
    add_text(s, Inches(0.75), Inches(11.8), Inches(6.0), Inches(0.4),
             "The 14th — gone before you opened your phone.",
             size=12, color=SLATE_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE FIX
# ────────────────────────────────────────────────────────────────────────────


def slide_2_fix(prs):
    s = blank(prs)
    gradient_background(s, PURPLE_TINT, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 2)

    add_eyebrow(s, 1.5, "THE FIX", color=PURPLE)

    # Headline
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.5), Inches(3.2),
             "Add it to your\nFramedrops\ncalendar.",
             size=50, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # Subhead
    add_text(s, Inches(0.5), Inches(5.4), Inches(6.5), Inches(0.6),
             "Every event reminder lands in your inbox — before the shoot.",
             size=16, color=SLATE_500, line_spacing=1.4,
             align=PP_ALIGN.LEFT)

    # Calendar mock with FOUR scheduled events in different accent colors
    draw_calendar_mock(s, 0.75, 6.4, 6.0, 3.5,
                       header_color=PURPLE,
                       events={
                           9:  INDIGO,
                           14: EMERALD,
                           22: AMBER,
                           27: PINK,
                       })

    # Below it — an email notification mock to make the "email reminder"
    # benefit visual and concrete.
    draw_email_notification(s, 0.75, 10.15, 6.0, 1.7)

    add_footer(s)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE CTA
# ────────────────────────────────────────────────────────────────────────────


def slide_3_cta(prs):
    s = blank(prs)
    gradient_background(s, INDIGO_TINT, WHITE)

    add_brand_mark(s)
    add_step_indicator(s, 3)

    add_eyebrow(s, 1.5, "TRY IT", color=PURPLE)

    # Headline
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.5), Inches(2.2),
             "Never miss a\nshoot again.",
             size=52, color=INK, bold=True, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # Subhead
    add_text(s, Inches(0.5), Inches(4.45), Inches(6.5), Inches(0.6),
             "Calendar + email reminders, built for photographers.",
             size=16, color=SLATE_500, line_spacing=1.4,
             align=PP_ALIGN.LEFT)

    # Phone mock — vertical, centered
    draw_phone_mock(s, x_in=2.05, y_in=5.4, w_in=3.4, h_in=4.5)

    # Offer card under the phone
    card_y = 10.2
    add_card(s, Inches(0.75), Inches(card_y), Inches(6.0), Inches(1.0),
             WHITE, corner=0.05, border=PURPLE_LINE, border_width=1.0)
    add_text(s, Inches(0.95), Inches(card_y + 0.18),
             Inches(5.7), Inches(0.4),
             "First client free",
             size=18, color=INK, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(0.95), Inches(card_y + 0.58),
             Inches(5.7), Inches(0.35),
             "3,000 photos · 30 days · No credit card",
             size=12, color=SLATE_500,
             align=PP_ALIGN.LEFT)

    # Big CTA pill
    add_link_pill(s, top_in=11.4, label="framedrops.in")


# ────────────────────────────────────────────────────────────────────────────


def build() -> None:
    prs = make_presentation()
    slide_1_pain(prs)
    slide_2_fix(prs)
    slide_3_cta(prs)
    prs.save(OUT_FILE)
    print(f"✓ Wrote {OUT_FILE}")
    print(f"  3 slides @ 9:16 vertical — English only, modern palette,")
    print(f"  with calendar / email / phone mocks. Link on every slide.")


if __name__ == "__main__":
    build()
