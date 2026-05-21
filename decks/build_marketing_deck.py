"""
Framedrops Marketing Deck — VIBRANT mobile edition.

Produces FrameDrops_Marketing_Deck.pptx — a 10-slide promotional deck
optimized for WhatsApp sharing (9:16 portrait). Bold, saturated, gradient-
heavy aesthetic — designed to make people stop scrolling.

Run:
    python3 build_marketing_deck.py

The output file lands in the same folder.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ─── VIBRANT palette — saturated, neon-leaning, not corporate-safe ────────
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_DEEP = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_NEON = RGBColor(0xA8, 0x55, 0xF7)
VIOLET = RGBColor(0x4C, 0x1D, 0x95)
INDIGO = RGBColor(0x3B, 0x0D, 0x86)
INDIGO_DEEP = RGBColor(0x1E, 0x1B, 0x4B)
PINK_HOT = RGBColor(0xEC, 0x48, 0x99)
PINK_NEON = RGBColor(0xF4, 0x3F, 0x5E)
ORANGE_HOT = RGBColor(0xFB, 0x7B, 0x2C)
CORAL = RGBColor(0xFF, 0x6B, 0x6B)
SUNSET = RGBColor(0xF9, 0x71, 0x6B)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CYAN_BRIGHT = RGBColor(0x22, 0xD3, 0xEE)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GREEN_NEON = RGBColor(0x34, 0xD3, 0x99)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
GOLD = RGBColor(0xFD, 0xE0, 0x47)
MAGENTA = RGBColor(0xD9, 0x46, 0xEF)

NAVY_BLACK = RGBColor(0x0B, 0x10, 0x26)
NAVY = RGBColor(0x14, 0x1B, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFD, 0xF4, 0xFF)
INK = RGBColor(0x1A, 0x0B, 0x3F)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT_SLATE = RGBColor(0xA1, 0xA1, 0xAA)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x10, 0xB9, 0x81)

# 9:16 portrait — WhatsApp-friendly.
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.33)

OUT_FILE = Path(__file__).parent / "FrameDrops_Marketing_Deck.pptx"
WEBSITE = "https://framedrops.in"


# ─── Core helpers ─────────────────────────────────────────────────────────


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def apply_gradient_fill(shape, stops, angle: float = 90) -> None:
    """Apply a real OOXML gradient fill to a shape.

    stops: list of (position_0_to_100, RGBColor) tuples.
    angle: 0 = left-to-right, 90 = top-to-bottom, 45 = diagonal.
    """
    sp_pr = shape.fill._xPr.spPr if hasattr(shape.fill._xPr, 'spPr') else shape._element.spPr
    # Remove any existing fill child
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill'):
        for el in sp_pr.findall(qn(tag)):
            sp_pr.remove(el)

    grad = etree.SubElement(sp_pr, qn('a:gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')

    gs_lst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, color in stops:
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', _rgb_hex(color))

    lin = etree.SubElement(grad, qn('a:lin'))
    # OOXML angle is in 60000ths of a degree; 0 = horizontal left-to-right.
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '0')

    # Make sure the spPr has the gradient where solid fill used to be.
    # python-pptx's shape.fill.* will now reflect "non-standard" — but it renders fine.


def add_gradient_rect(slide, left, top, width, height, stops, angle=90,
                       corner: float = 0.0):
    """Rounded or square rectangle with gradient fill."""
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if corner > 0 else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, left, top, width, height)
    if corner > 0:
        s.adjustments[0] = corner
    s.line.fill.background()
    apply_gradient_fill(s, stops, angle)
    return s


def add_gradient_oval(slide, left, top, size, stops, angle=45):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    o.line.fill.background()
    apply_gradient_fill(o, stops, angle)
    return o


def add_solid_oval(slide, left, top, size, color):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    o.line.fill.background()
    o.fill.solid()
    o.fill.fore_color.rgb = color
    return o


def gradient_background(slide, stops, angle: float = 135) -> None:
    """Fill the whole slide with a gradient."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    apply_gradient_fill(bg, stops, angle)


def fill_solid(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_text(
    slide, left, top, width, height, text: str, *,
    size: int = 24, color: RGBColor = WHITE, bold: bool = False,
    align: int = PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Inter"
    return tb


def add_gradient_pill(slide, left, top, width, height, text: str,
                       stops, fg: RGBColor = WHITE, size: int = 12, angle=0):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    apply_gradient_fill(pill, stops, angle)
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Inter"
    return pill


def add_pill(slide, left, top, width, height, text: str, fill: RGBColor,
             fg: RGBColor = WHITE, size: int = 12):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Inter"
    return pill


def add_card(slide, left, top, width, height, fill: RGBColor,
             border: RGBColor = None, border_width: float = 1.5,
             corner: float = 0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card


def add_gradient_card(slide, left, top, width, height, stops,
                       angle=135, corner: float = 0.06,
                       border: RGBColor = None, border_width: float = 1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.adjustments[0] = corner
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    apply_gradient_fill(card, stops, angle)
    return card


def add_strike_price(slide, left, top, width, text: str, size: int = 18,
                      color: RGBColor = LIGHT_SLATE):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Inter"
    rPr = run._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")
    return tb


# ────────────────────────────────────────────────────────────────────────────
# Slides
# ────────────────────────────────────────────────────────────────────────────


def slide_1_hook(prs):
    s = blank(prs)
    # Deep cosmic gradient: indigo → violet → magenta
    gradient_background(s, [
        (0, INDIGO_DEEP),
        (50, VIOLET),
        (100, PINK_HOT),
    ], angle=135)

    # Big bright glow orbs
    add_gradient_oval(s, Inches(-2), Inches(7), Inches(9), [
        (0, MAGENTA), (100, ORANGE_HOT),
    ], angle=45)
    add_gradient_oval(s, Inches(4), Inches(-2), Inches(7), [
        (0, CYAN_BRIGHT), (100, PURPLE_NEON),
    ], angle=315)

    # Gradient brand pill at top
    add_gradient_pill(s, Inches(2.5), Inches(0.8), Inches(2.5), Inches(0.55),
                      "✨ FRAMEDROPS", [
                          (0, YELLOW), (100, ORANGE_HOT)
                      ], fg=INK, size=13)

    # Headline — huge & punchy
    add_text(s, Inches(0.5), Inches(2.5), Inches(6.5), Inches(4),
             "Tired of\nGoogle Drive\n+ paper lists?",
             size=62, color=WHITE, bold=True, line_spacing=1.0)

    # Yellow accent underline behind a word for pop
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(6.5),
                                     Inches(3.5), Inches(0.18))
    apply_gradient_fill(underline, [(0, YELLOW), (100, ORANGE_HOT)], angle=0)
    underline.line.fill.background()

    add_text(s, Inches(0.5), Inches(7.0), Inches(6.5), Inches(2.5),
             "The proper selection workflow\nfor Indian photographers.",
             size=22, color=WHITE, line_spacing=1.3)

    # Bottom: 3 trust chips
    chips = [
        ("🇮🇳 Made in India", [(0, ORANGE_HOT), (100, PINK_NEON)]),
        ("⚡ 2-min signup", [(0, CYAN_BRIGHT), (100, PURPLE)]),
        ("🎁 300 free", [(0, GREEN_NEON), (100, TEAL)]),
    ]
    x = 0.5
    for text, stops in chips:
        add_gradient_pill(s, Inches(x), Inches(12.0), Inches(2.2), Inches(0.55),
                          text, stops, fg=WHITE, size=11, angle=45)
        x += 2.35

    # Decorative sparkles
    for sx, sy, sz in [(5.5, 3.2, 0.18), (1.5, 8.5, 0.14),
                        (6.2, 9.5, 0.22), (0.8, 11.0, 0.12)]:
        spark = s.shapes.add_shape(MSO_SHAPE.STAR_4_POINT,
                                     Inches(sx), Inches(sy),
                                     Inches(sz), Inches(sz))
        spark.fill.solid()
        spark.fill.fore_color.rgb = YELLOW
        spark.line.fill.background()


def slide_2_pain(prs):
    s = blank(prs)
    # Warm sunset background: orange → coral → pink
    gradient_background(s, [
        (0, RGBColor(0xFE, 0xF3, 0xC7)),
        (50, RGBColor(0xFC, 0xD3, 0x4D)),
        (100, ORANGE_HOT),
    ], angle=135)

    # Decorative dark blob top-right
    add_gradient_oval(s, Inches(4.5), Inches(-2), Inches(6), [
        (0, PINK_HOT), (100, CORAL),
    ], angle=180)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(2.5), Inches(0.55),
                      "🔥 THE PAIN", [(0, RED), (100, PINK_NEON)], size=13)

    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(2),
             "Sound familiar?", size=56, color=INK, bold=True,
             line_spacing=1.0)

    pains = [
        ("📁", "Upload 1,500 photos to Drive",
         "30+ min waiting. Storage filling up.",
         [(0, PURPLE), (100, PINK_HOT)]),
        ("🔗", "Share link on WhatsApp",
         "Client pinch-zooms on tiny phone screen.",
         [(0, CYAN), (100, PURPLE)]),
        ("📝", "Names written on paper",
         "“IMG_2845, IMG_2912, IMG_3001…”",
         [(0, ORANGE_HOT), (100, PINK_NEON)]),
        ("💬", "WhatsApp the paper photo",
         "Or typed list. Or both. Or neither.",
         [(0, MAGENTA), (100, PURPLE_DEEP)]),
        ("🔍", "Decode the handwriting",
         "Search file-by-file. Miss one. Complain.",
         [(0, TEAL), (100, INDIGO)]),
        ("😩", "Repeat. Every. Wedding.",
         "Forever and ever and ever.",
         [(0, PINK_NEON), (100, VIOLET)]),
    ]

    y = 3.7
    for emoji, head, sub, stops in pains:
        # Gradient circle for icon
        add_gradient_oval(s, Inches(0.5), Inches(y), Inches(0.9), stops, angle=45)
        add_text(s, Inches(0.5), Inches(y + 0.1), Inches(0.9), Inches(0.7),
                 emoji, size=24, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Text block
        add_text(s, Inches(1.6), Inches(y + 0.05), Inches(5.7), Inches(0.5),
                 head, size=17, color=INK, bold=True, line_spacing=1.1)
        add_text(s, Inches(1.6), Inches(y + 0.5), Inches(5.7), Inches(0.5),
                 sub, size=12, color=SLATE, line_spacing=1.3)
        y += 1.45


def slide_3_compare(prs):
    s = blank(prs)
    # Soft cool gradient: lavender → mint
    gradient_background(s, [
        (0, RGBColor(0xF5, 0xE6, 0xFF)),
        (100, RGBColor(0xD8, 0xF5, 0xEC)),
    ], angle=135)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(2.5), Inches(0.55),
                      "✨ THE FIX", [(0, GREEN_NEON), (100, TEAL)], size=13)
    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(2),
             "There's a\nbetter way.",
             size=56, color=INK, bold=True, line_spacing=1.0)

    rows = [
        ("Spam screenshots", "One branded gallery"),
        ("Lost selections", "Auto-saved client picks"),
        ("Manual file copy", "Smart local transfer"),
        ("Looks like a freelancer", "Looks like a studio"),
        ("7-day back-and-forth", "Selections in 48 hours"),
        ("Free but chaotic", "₹149 saves you 6 hours"),
    ]

    # Column headers as gradient bars
    add_gradient_card(s, Inches(0.4), Inches(4.5), Inches(3.3), Inches(0.7),
                      [(0, RED), (100, PINK_NEON)], angle=90, corner=0.15)
    add_text(s, Inches(0.4), Inches(4.65), Inches(3.3), Inches(0.4),
             "❌ OLD WAY", size=15, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)

    add_gradient_card(s, Inches(3.85), Inches(4.5), Inches(3.3), Inches(0.7),
                      [(0, GREEN_NEON), (100, TEAL)], angle=90, corner=0.15)
    add_text(s, Inches(3.85), Inches(4.65), Inches(3.3), Inches(0.4),
             "✅ FRAMEDROPS", size=15, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)

    y = 5.4
    for bad, good in rows:
        # Bad cell — soft red gradient card
        add_gradient_card(s, Inches(0.4), Inches(y), Inches(3.3), Inches(1.1),
                          [(0, RGBColor(0xFF, 0xE4, 0xE6)),
                           (100, RGBColor(0xFE, 0xCA, 0xCA))],
                          angle=135, corner=0.12)
        add_text(s, Inches(0.55), Inches(y + 0.25), Inches(3.0), Inches(0.7),
                 bad, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.2)
        # Good cell — vibrant green gradient card
        add_gradient_card(s, Inches(3.85), Inches(y), Inches(3.3), Inches(1.1),
                          [(0, PURPLE), (100, CYAN)],
                          angle=135, corner=0.12)
        add_text(s, Inches(4.0), Inches(y + 0.25), Inches(3.0), Inches(0.7),
                 good, size=13, color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        y += 1.2


def slide_4_how(prs):
    s = blank(prs)
    # Deep midnight → electric purple gradient
    gradient_background(s, [
        (0, NAVY_BLACK),
        (50, INDIGO_DEEP),
        (100, PURPLE_DEEP),
    ], angle=180)

    # Glow accent
    add_gradient_oval(s, Inches(-2), Inches(9), Inches(8), [
        (0, MAGENTA), (100, PURPLE_DEEP),
    ], angle=45)
    add_gradient_oval(s, Inches(5), Inches(-1), Inches(5), [
        (0, CYAN_BRIGHT), (100, PURPLE),
    ], angle=180)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(3), Inches(0.55),
                      "⚡ HOW IT WORKS", [(0, YELLOW), (100, ORANGE_HOT)],
                      fg=INK, size=13)

    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(3),
             "Shoot to\nselection in\n3 steps.",
             size=54, color=WHITE, bold=True, line_spacing=1.0)

    steps = [
        ("1", "Upload",
         "Drag-drop your shoot. Auto-compressed.\n25 MB → 250 KB.",
         "📤",
         [(0, CYAN_BRIGHT), (100, PURPLE_NEON)]),
        ("2", "Share",
         "One branded gallery link via WhatsApp.\nClient taps to view.",
         "🔗",
         [(0, PURPLE_NEON), (100, PINK_HOT)]),
        ("3", "Get selections",
         "Client taps the heart on favorites.\nNo paper. No decoding.",
         "❤️",
         [(0, PINK_NEON), (100, ORANGE_HOT)]),
    ]

    y = 5.5
    for num, title, body, emoji, stops in steps:
        # Numbered gradient circle
        add_gradient_oval(s, Inches(0.5), Inches(y), Inches(1.0), stops, angle=135)
        add_text(s, Inches(0.5), Inches(y + 0.12), Inches(1.0), Inches(0.8),
                 num, size=30, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title + emoji
        add_text(s, Inches(1.7), Inches(y - 0.05), Inches(5.5), Inches(0.6),
                 f"{emoji}  {title}", size=24, color=WHITE, bold=True)
        # Body
        add_text(s, Inches(1.7), Inches(y + 0.6), Inches(5.5), Inches(1.5),
                 body, size=13, color=LIGHT_SLATE, line_spacing=1.4)
        y += 2.4


def slide_5_india(prs):
    s = blank(prs)
    # India-flavored gradient: saffron → white-pink → green-teal
    gradient_background(s, [
        (0, ORANGE_HOT),
        (50, PINK_NEON),
        (100, VIOLET),
    ], angle=160)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(3.5), Inches(0.55),
                      "🇮🇳  MADE IN INDIA", [(0, GOLD), (100, YELLOW)],
                      fg=INK, size=13)

    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(3),
             "Made for the\nIndian wedding.",
             size=48, color=WHITE, bold=True, line_spacing=1.0)

    perks = [
        ("🗣️", "Three Indian languages", "English · Telugu · Hindi",
         [(0, ORANGE_HOT), (100, RED)]),
        ("💳", "INR + GST native", "Razorpay UPI checkout",
         [(0, PURPLE), (100, MAGENTA)]),
        ("📱", "Mobile-first client view", "Your clients live on phones",
         [(0, CYAN), (100, TEAL)]),
        ("🏦", "Wallet → any Indian bank", "NEFT payouts",
         [(0, GREEN_NEON), (100, TEAL)]),
        ("💸", "No monthly subscription", "Pay only when you publish",
         [(0, YELLOW), (100, ORANGE_HOT)]),
    ]

    y = 5.0
    for emoji, head, sub, stops in perks:
        # Gradient icon orb
        add_gradient_oval(s, Inches(0.5), Inches(y), Inches(0.9), stops, angle=45)
        add_text(s, Inches(0.5), Inches(y + 0.1), Inches(0.9), Inches(0.7),
                 emoji, size=24, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.6), Inches(y + 0.05), Inches(5.7), Inches(0.5),
                 head, size=17, color=WHITE, bold=True, line_spacing=1.1)
        add_text(s, Inches(1.6), Inches(y + 0.5), Inches(5.7), Inches(0.5),
                 sub, size=12, color=RGBColor(0xFE, 0xE0, 0xF5), line_spacing=1.3)
        y += 1.4


def slide_6_pricing(prs):
    s = blank(prs)
    # Light dreamy gradient: pink → lavender → cyan
    gradient_background(s, [
        (0, RGBColor(0xFD, 0xE6, 0xFF)),
        (50, RGBColor(0xE9, 0xD5, 0xFF)),
        (100, RGBColor(0xCF, 0xFB, 0xFD)),
    ], angle=160)

    # Launch banner — vibrant gradient
    add_gradient_card(s, Inches(0.4), Inches(0.6), Inches(6.7), Inches(1.0),
                      [(0, PURPLE_DEEP), (50, PINK_HOT), (100, ORANGE_HOT)],
                      angle=90, corner=0.1)
    add_gradient_pill(s, Inches(0.6), Inches(0.75), Inches(2.0), Inches(0.4),
                      "🚀 LAUNCH OFFER",
                      [(0, YELLOW), (100, ORANGE_HOT)], fg=INK, size=10)
    add_text(s, Inches(0.6), Inches(1.2), Inches(6.3), Inches(0.4),
             "First 100 photographers · Locked-in rates",
             size=12, color=WHITE, bold=True)

    add_text(s, Inches(0.4), Inches(1.9), Inches(6.7), Inches(1.2),
             "Simple,\nhonest pricing.",
             size=42, color=INK, bold=True, line_spacing=1.0)

    # Free tier strip — bold gradient
    add_gradient_card(s, Inches(0.4), Inches(4.5), Inches(6.7), Inches(1.0),
                      [(0, INDIGO_DEEP), (100, PURPLE_DEEP)],
                      angle=90, corner=0.1)
    add_text(s, Inches(0.6), Inches(4.6), Inches(6.5), Inches(0.5),
             "🎁  First 300 images — FREE for life",
             size=15, color=WHITE, bold=True)
    add_text(s, Inches(0.6), Inches(5.0), Inches(6.5), Inches(0.4),
             "No credit card. Setup in 2 minutes.",
             size=11, color=RGBColor(0xC8, 0xB8, 0xFF))

    # Three tier cards — bold gradient backgrounds
    tiers = [
        ("STARTER", "149", "199", "Up to 1,000 photos", "₹0.15/img",
         [(0, CYAN_BRIGHT), (100, PURPLE)], False),
        ("⭐ MOST POPULAR", "229", "299", "Up to 2,000 photos", "₹0.11/img",
         [(0, PINK_HOT), (50, MAGENTA), (100, PURPLE)], True),
        ("WEDDING PRO", "299", "399", "Up to 3,000 photos", "₹0.10/img",
         [(0, ORANGE_HOT), (100, PINK_NEON)], False),
    ]

    y = 5.9
    for name, price, was, range_text, per_img, stops, popular in tiers:
        h = 2.05 if popular else 1.9
        add_gradient_card(s, Inches(0.4), Inches(y), Inches(6.7), Inches(h),
                          stops, angle=135, corner=0.07,
                          border=YELLOW if popular else None,
                          border_width=3 if popular else 0)
        # Tier name
        add_text(s, Inches(0.7), Inches(y + 0.2), Inches(3.5), Inches(0.4),
                 name, size=11, color=WHITE, bold=True)
        # Price block
        add_text(s, Inches(0.7), Inches(y + 0.6), Inches(2.5), Inches(1.0),
                 f"₹{price}", size=44, color=WHITE, bold=True,
                 line_spacing=1.0)
        # Strike — yellowish so it pops on dark gradient
        add_strike_price(s, Inches(0.7), Inches(y + 1.5), Inches(2.0),
                          f"₹{was}", size=14,
                          color=RGBColor(0xFF, 0xE6, 0x9C))
        # Right side: range + per image
        add_text(s, Inches(3.5), Inches(y + 0.6), Inches(3.4), Inches(0.5),
                 range_text, size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.RIGHT)
        add_text(s, Inches(3.5), Inches(y + 1.05), Inches(3.4), Inches(0.5),
                 per_img, size=13, color=YELLOW, bold=True, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(3.5), Inches(y + 1.55), Inches(3.4), Inches(0.4),
                 "Pay only when you publish", size=10,
                 color=RGBColor(0xF0, 0xE5, 0xFF), align=PP_ALIGN.RIGHT)
        y += h + 0.15


def slide_7_premium(prs):
    s = blank(prs)
    # Gold/luxury gradient
    gradient_background(s, [
        (0, INDIGO_DEEP),
        (50, PURPLE_DEEP),
        (100, RGBColor(0x6B, 0x21, 0xA8)),
    ], angle=180)

    # Gold glow
    add_gradient_oval(s, Inches(-2), Inches(6), Inches(8), [
        (0, YELLOW), (100, ORANGE_HOT),
    ], angle=45)
    add_gradient_oval(s, Inches(5), Inches(-2), Inches(6), [
        (0, MAGENTA), (100, PINK_HOT),
    ], angle=180)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(3), Inches(0.55),
                      "👑 LOOK PREMIUM", [(0, GOLD), (100, ORANGE_HOT)],
                      fg=INK, size=13)

    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(4),
             "Look like a\n₹2-lakh\nphotographer.",
             size=58, color=WHITE, bold=True, line_spacing=1.0)

    # Yellow accent underline
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(5.4),
                                     Inches(2.5), Inches(0.15))
    apply_gradient_fill(underline, [(0, YELLOW), (100, ORANGE_HOT)], angle=0)
    underline.line.fill.background()

    perks = [
        "Custom studio branding on every gallery",
        "Watermarked, screenshot-protected previews",
        "Anti-download protection on originals",
        "Clients see YOUR name — not 'Framedrops'",
        "Look professional from day one",
    ]

    y = 6.5
    for line in perks:
        # Gradient check icon
        add_gradient_oval(s, Inches(0.5), Inches(y), Inches(0.45),
                          [(0, GREEN_NEON), (100, CYAN_BRIGHT)], angle=45)
        add_text(s, Inches(0.5), Inches(y + 0.02), Inches(0.45), Inches(0.45),
                 "✓", size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.1), Inches(y + 0.03), Inches(6.0), Inches(0.5),
                 line, size=15, color=WHITE, line_spacing=1.3)
        y += 0.95


def slide_8_value(prs):
    s = blank(prs)
    # Vibrant green-cyan-purple gradient
    gradient_background(s, [
        (0, TEAL),
        (50, CYAN),
        (100, PURPLE_DEEP),
    ], angle=160)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(3), Inches(0.55),
                      "💰 THE NUMBERS", [(0, GREEN_NEON), (100, YELLOW)],
                      fg=INK, size=13)

    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(3.5),
             "₹149 pays\nfor itself in\n15 minutes.",
             size=48, color=WHITE, bold=True, line_spacing=1.0)

    stats = [
        ("6+ hrs", "Saved per wedding — no decoding",
         [(0, ORANGE_HOT), (100, PINK_HOT)]),
        ("Zero", "File-by-file Drive searching",
         [(0, YELLOW), (100, ORANGE_HOT)]),
        ("48 hrs", "Shoot → selection (vs 7+ days)",
         [(0, CYAN_BRIGHT), (100, PURPLE)]),
        ("₹500-2k/hr", "What your time is worth",
         [(0, GREEN_NEON), (100, TEAL)]),
    ]

    y = 6.0
    for big, sub, stops in stats:
        add_gradient_card(s, Inches(0.5), Inches(y), Inches(6.5), Inches(1.45),
                          stops, angle=135, corner=0.1)
        add_text(s, Inches(0.85), Inches(y + 0.25), Inches(3.0), Inches(1.0),
                 big, size=30, color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Divider
        div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(3.8), Inches(y + 0.45),
                                   Inches(0.04), Inches(0.55))
        div.fill.solid()
        div.fill.fore_color.rgb = WHITE
        div.line.fill.background()
        add_text(s, Inches(4.0), Inches(y + 0.3), Inches(2.9), Inches(0.9),
                 sub, size=12, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.35)
        y += 1.6


def slide_9_vs(prs):
    s = blank(prs)
    # Bold violet gradient with bright accent
    gradient_background(s, [
        (0, RGBColor(0xFD, 0xF4, 0xFF)),
        (100, RGBColor(0xE9, 0xD5, 0xFF)),
    ], angle=160)

    add_gradient_pill(s, Inches(0.5), Inches(0.8), Inches(3), Inches(0.55),
                      "🆚 COMPARISON", [(0, MAGENTA), (100, PURPLE_DEEP)],
                      size=13)
    add_text(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(2.5),
             "Why not free\ncloud + paper?",
             size=46, color=INK, bold=True, line_spacing=1.0)
    add_text(s, Inches(0.5), Inches(4.0), Inches(6.5), Inches(0.6),
             "Because 'free' has a hidden cost.",
             size=15, color=SLATE)

    # Old way card — washed sunset
    add_gradient_card(s, Inches(0.4), Inches(5.0), Inches(6.7), Inches(3.5),
                      [(0, RGBColor(0xFF, 0xE4, 0xE6)),
                       (100, RGBColor(0xFC, 0xA5, 0xA5))],
                      angle=135, corner=0.07, border=RED, border_width=2)
    add_text(s, Inches(0.7), Inches(5.2), Inches(6.0), Inches(0.4),
             "THE OLD WAY", size=11, color=RED, bold=True)
    add_text(s, Inches(0.7), Inches(5.6), Inches(6.0), Inches(0.6),
             "Free cloud + paper lists",
             size=22, color=INK, bold=True)
    add_text(s, Inches(0.7), Inches(6.3), Inches(6.0), Inches(2.0),
             "❌  Clients write file names on paper\n"
             "❌  You decode handwriting, file-by-file\n"
             "❌  No branding — looks like a hobbyist\n"
             "❌  Storage nags for paid upgrade\n"
             "❌  Misread one name → angry client",
             size=12, color=INK, line_spacing=1.7)

    # Framedrops card — vibrant
    add_gradient_card(s, Inches(0.4), Inches(8.7), Inches(6.7), Inches(4.0),
                      [(0, PURPLE_DEEP), (50, PINK_HOT), (100, ORANGE_HOT)],
                      angle=135, corner=0.07, border=YELLOW, border_width=3)
    add_gradient_pill(s, Inches(0.7), Inches(8.9), Inches(2.0), Inches(0.4),
                      "⭐ FRAMEDROPS", [(0, YELLOW), (100, ORANGE_HOT)],
                      fg=INK, size=10)
    add_text(s, Inches(0.7), Inches(9.4), Inches(6.0), Inches(0.6),
             "₹149 per album",
             size=32, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(10.05), Inches(6.0), Inches(0.5),
             "(only when you publish)",
             size=12, color=RGBColor(0xFF, 0xE9, 0xF7))
    add_text(s, Inches(0.7), Inches(10.6), Inches(6.0), Inches(2.2),
             "✅  Pay zero when you don't shoot\n"
             "✅  INR-native — no FX surprises\n"
             "✅  English, Telugu, Hindi\n"
             "✅  GST invoices + Razorpay UPI\n"
             "✅  Built in India, for Indian weddings",
             size=12, color=WHITE, line_spacing=1.7)


def slide_10_cta(prs):
    s = blank(prs)
    # Epic finale gradient
    gradient_background(s, [
        (0, INDIGO_DEEP),
        (30, VIOLET),
        (70, PINK_HOT),
        (100, ORANGE_HOT),
    ], angle=160)

    # Big glow
    add_gradient_oval(s, Inches(-3), Inches(8), Inches(12), [
        (0, MAGENTA), (100, PURPLE_DEEP),
    ], angle=45)
    add_gradient_oval(s, Inches(5), Inches(-2), Inches(6), [
        (0, YELLOW), (100, ORANGE_HOT),
    ], angle=180)

    add_gradient_pill(s, Inches(2.0), Inches(1.0), Inches(3.5), Inches(0.6),
                      "✨ START FREE TODAY",
                      [(0, YELLOW), (100, ORANGE_HOT)], fg=INK, size=14)

    add_text(s, Inches(0.5), Inches(2.5), Inches(6.5), Inches(5),
             "Your next\nwedding\ndeserves\nbetter.",
             size=64, color=WHITE, bold=True, line_spacing=1.0)

    # Trust line
    add_text(s, Inches(0.5), Inches(7.8), Inches(6.5), Inches(0.6),
             "🎁 300 free  ·  💳 No card  ·  ⚡ 2 min",
             size=16, color=RGBColor(0xFF, 0xE6, 0xF8),
             align=PP_ALIGN.CENTER)

    # Big CTA URL card with shadow effect
    # Shadow card (back)
    add_gradient_card(s, Inches(0.6), Inches(9.0), Inches(6.3), Inches(2.0),
                      [(0, PURPLE_DEEP), (100, PINK_HOT)],
                      angle=135, corner=0.08)
    # Front white card
    add_card(s, Inches(0.5), Inches(8.85), Inches(6.3), Inches(2.0),
             WHITE, corner=0.08)

    add_text(s, Inches(0.5), Inches(9.0), Inches(6.3), Inches(0.5),
             "TRY FRAMEDROPS FREE", size=12, color=PURPLE_DEEP, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(9.5), Inches(6.3), Inches(0.9),
             WEBSITE, size=26, color=INK, bold=True, align=PP_ALIGN.CENTER)
    # Gradient underline
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(2.5), Inches(10.4),
                                     Inches(2.3), Inches(0.06))
    apply_gradient_fill(underline, [(0, PINK_HOT), (100, ORANGE_HOT)], angle=0)
    underline.line.fill.background()

    # Tagline
    add_text(s, Inches(0.5), Inches(11.5), Inches(6.5), Inches(0.6),
             "Stop chasing.  Start delivering.",
             size=20, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────


def build() -> None:
    prs = make_presentation()
    slide_1_hook(prs)
    slide_2_pain(prs)
    slide_3_compare(prs)
    slide_4_how(prs)
    slide_5_india(prs)
    slide_6_pricing(prs)
    slide_7_premium(prs)
    slide_8_value(prs)
    slide_9_vs(prs)
    slide_10_cta(prs)
    prs.save(OUT_FILE)
    print(f"✓ Wrote {OUT_FILE}")


if __name__ == "__main__":
    build()
