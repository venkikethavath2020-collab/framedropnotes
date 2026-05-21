"""
Framedrops Marketing Deck — Desktop edition.

Produces FrameDrops_Marketing_Deck_Desktop.pptx — a 10-slide promotional deck
optimized for desktop/projector viewing (16:9 widescreen, multi-column layouts,
richer visual treatment).

Same content as the mobile deck (build_marketing_deck.py), but laid out for
landscape with larger typography, layered backgrounds, and decorative accents.

Run:
    python3 build_marketing_deck_desktop.py

The output file lands in the same folder.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Brand palette (matches framedrops/src/plugins/vuetify.ts) ────────────
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
DEEP_PURPLE = RGBColor(0x4F, 0x46, 0xE5)
INDIGO = RGBColor(0x31, 0x2E, 0x81)
PINK = RGBColor(0xEC, 0x48, 0x99)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
MARIGOLD = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE_LAVENDER = RGBColor(0xFA, 0xF5, 0xFF)

# 16:9 widescreen — desktop / projector.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_FILE = Path(__file__).parent / "FrameDrops_Marketing_Deck_Desktop.pptx"
WEBSITE = "https://framedrops.in"


# ─── Helpers ──────────────────────────────────────────────────────────────


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_glow(slide, left, top, size, color: RGBColor):
    """Decorative blurred-look circle (no real blur, but large soft purple disc)."""
    g = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    g.fill.solid()
    g.fill.fore_color.rgb = color
    g.line.fill.background()
    return g


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 24,
    color: RGBColor = WHITE,
    bold: bool = False,
    align: int = PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Inter"
    return tb


def add_pill(slide, left, top, width, height, text: str, fill: RGBColor,
             fg: RGBColor = WHITE, size: int = 12):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    tf = pill.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Inter"
    return pill


def add_card(slide, left, top, width, height, fill: RGBColor,
             border: RGBColor = None, border_width: float = 1.5,
             corner: float = 0.05):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card


def add_accent_bar(slide, left, top, height, color: RGBColor, width=Inches(0.06)):
    """Vertical accent bar — a brand-recognizable thin strip on the side of a section."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_strike_price(slide, left, top, width, text: str, size: int = 20,
                      color: RGBColor = SLATE_400):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Inter"
    rPr = run._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")
    return tb


def add_side_rail(slide, color: RGBColor = PURPLE):
    """Vertical brand strip down the left edge — present on every slide for consistency."""
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, 0, Inches(0.12), SLIDE_H)
    rail.fill.solid()
    rail.fill.fore_color.rgb = color
    rail.line.fill.background()
    return rail


def add_page_number(slide, num: int, total: int, on_dark: bool = False):
    color = SLATE_400 if not on_dark else SLATE_300
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
                f"{num:02d} / {total:02d}", size=10, color=color,
                bold=True, align=PP_ALIGN.RIGHT)


def add_brand_lockup(slide, on_dark: bool = False):
    color = PURPLE if not on_dark else WHITE
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(3), Inches(0.4),
                "FRAMEDROPS", size=11, color=color, bold=True, align=PP_ALIGN.LEFT)
    # tiny dot accent
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(1.8), Inches(0.55), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PURPLE if not on_dark else MARIGOLD
    dot.line.fill.background()


# ────────────────────────────────────────────────────────────────────────────
# Slides
# ────────────────────────────────────────────────────────────────────────────


def slide_1_hook(prs):
    s = blank(prs)
    fill_background(s, NAVY)

    # Layered gradient — three overlapping circles in indigo/purple/pink for depth.
    add_glow(s, Inches(-3), Inches(-3), Inches(9), INDIGO)
    add_glow(s, Inches(7), Inches(-2), Inches(8), DEEP_PURPLE)
    add_glow(s, Inches(9), Inches(4), Inches(6), PURPLE)

    # Diagonal accent stripe — bottom right
    stripe = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                  Inches(9), Inches(6.2), Inches(5), Inches(0.4))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = MARIGOLD
    stripe.line.fill.background()

    add_brand_lockup(s, on_dark=True)

    # Left content
    add_pill(s, Inches(0.5), Inches(2.0), Inches(2.5), Inches(0.5),
             "🇮🇳  MADE IN INDIA", PURPLE, WHITE, size=11)

    add_textbox(s, Inches(0.5), Inches(2.8), Inches(7.5), Inches(3),
                "Tired of\nGoogle Drive\n+ paper lists?",
                size=72, color=WHITE, bold=True, line_spacing=1.0)

    add_textbox(s, Inches(0.5), Inches(5.7), Inches(7.5), Inches(1.5),
                "The proper photo-selection workflow\nfor Indian photographers.",
                size=22, color=SLATE_400, line_spacing=1.3)

    # Right-side visual: stacked "card" preview suggesting a gallery
    card_x = Inches(9.0)
    card_y = Inches(2.0)
    # Outer card (back)
    back = add_card(s, card_x + Inches(0.3), card_y + Inches(0.3),
                    Inches(3.5), Inches(4.5), DEEP_PURPLE, corner=0.06)
    # Mid card
    mid = add_card(s, card_x + Inches(0.15), card_y + Inches(0.15),
                   Inches(3.5), Inches(4.5), PURPLE, corner=0.06)
    # Front card with mock gallery
    front = add_card(s, card_x, card_y, Inches(3.5), Inches(4.5), WHITE, corner=0.06)
    add_textbox(s, card_x + Inches(0.3), card_y + Inches(0.3),
                Inches(3), Inches(0.4),
                "STUDIO  •  WEDDING", size=9, color=PURPLE, bold=True)
    add_textbox(s, card_x + Inches(0.3), card_y + Inches(0.65),
                Inches(3), Inches(0.5),
                "Anjali + Rohan", size=18, color=NAVY, bold=True)
    # Grid of thumbnail squares
    for row in range(4):
        for col in range(3):
            tx = card_x + Inches(0.3) + Inches(col * 1.0)
            ty = card_y + Inches(1.4) + Inches(row * 0.7)
            thumb = add_card(s, tx, ty, Inches(0.9), Inches(0.6),
                             SLATE_200, corner=0.08)
            # Mark a couple as "favorited"
            if (row, col) in [(0, 1), (1, 0), (2, 2), (3, 1)]:
                # Replace with a purple-tinted thumb + small heart
                thumb.fill.fore_color.rgb = PURPLE_LAVENDER
                heart = s.shapes.add_shape(MSO_SHAPE.HEART,
                                             tx + Inches(0.65), ty + Inches(0.05),
                                             Inches(0.2), Inches(0.18))
                heart.fill.solid()
                heart.fill.fore_color.rgb = PURPLE
                heart.line.fill.background()

    add_page_number(s, 1, 10, on_dark=True)


def slide_2_pain(prs):
    s = blank(prs)
    fill_background(s, OFF_WHITE)
    add_side_rail(s, RED)
    add_brand_lockup(s)

    # Big section number watermark
    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "01", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.0), Inches(0.5),
             "THE PROBLEM", RED, WHITE, size=11)

    add_textbox(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(1.2),
                "Sound familiar?",
                size=54, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.9), Inches(7.5), Inches(0.6),
                "The hidden tax of every Indian wedding shoot.",
                size=16, color=SLATE_500)

    # 6 pain cards in a 3×2 grid on the right + emphasis text on left
    pains = [
        ("📁", "Upload 1,500 photos to Drive", "30+ min waiting. Storage fills up."),
        ("🔗", "Share link on WhatsApp", "Client pinch-zooms on phone."),
        ("📝", "Names on paper", "“IMG_2845, IMG_2912, IMG_3001…”"),
        ("💬", "WhatsApp the paper photo", "Or typed list. Or both."),
        ("🔍", "Decode the handwriting", "Search file-by-file. Miss one."),
        ("😩", "Repeat every wedding", "Forever."),
    ]
    grid_left = 0.5
    grid_top = 4.0
    cw, ch = 2.05, 1.5
    gap_x, gap_y = 0.15, 0.15
    for idx, (emoji, head, sub) in enumerate(pains):
        col = idx % 3
        row = idx // 3
        x = Inches(grid_left + col * (cw + gap_x))
        y = Inches(grid_top + row * (ch + gap_y))
        add_card(s, x, y, Inches(cw), Inches(ch), WHITE, border=SLATE_200, corner=0.08)
        add_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                    emoji, size=22, color=NAVY)
        add_textbox(s, x + Inches(0.2), y + Inches(0.65), Inches(cw - 0.4), Inches(0.45),
                    head, size=13, color=NAVY, bold=True, line_spacing=1.1)
        add_textbox(s, x + Inches(0.2), y + Inches(1.05), Inches(cw - 0.4), Inches(0.4),
                    sub, size=10, color=SLATE_500, line_spacing=1.2)

    # Right column emphasis
    right_x = Inches(7.4)
    add_card(s, right_x, Inches(4.0), Inches(5.3), Inches(3.15),
             NAVY, corner=0.04)
    add_textbox(s, right_x + Inches(0.4), Inches(4.3), Inches(4.5), Inches(0.5),
                "THE REAL COST", size=11, color=MARIGOLD, bold=True)
    add_textbox(s, right_x + Inches(0.4), Inches(4.7), Inches(4.5), Inches(2),
                "6+ hours\nper wedding",
                size=42, color=WHITE, bold=True, line_spacing=1.0)
    add_textbox(s, right_x + Inches(0.4), Inches(6.3), Inches(4.5), Inches(0.8),
                "Time you could spend\non your next shoot.",
                size=13, color=SLATE_400, line_spacing=1.3)

    add_page_number(s, 2, 10)


def slide_3_compare(prs):
    s = blank(prs)
    fill_background(s, WHITE)
    add_side_rail(s, GREEN)
    add_brand_lockup(s)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "02", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.2), Inches(0.5),
             "THE FIX", GREEN, WHITE, size=11)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(8), Inches(1.2),
                "There's a better way.",
                size=54, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.9), Inches(8), Inches(0.5),
                "Every painful step, replaced by one click.",
                size=16, color=SLATE_500)

    rows = [
        ("Upload to Google Drive",
         "Drag-drop — auto-compressed in browser"),
        ("Client opens Drive on phone",
         "Mobile-first branded gallery"),
        ("Names written on paper",
         "One-tap heart on favorites"),
        ("WhatsApp message with the list",
         "Selections auto-saved to dashboard"),
        ("You search file-by-file",
         "Smart local transfer — names match"),
        ("Looks like a freelancer",
         "Looks like a pro studio"),
    ]

    # Headers
    header_y = 4.0
    add_card(s, Inches(0.5), Inches(header_y), Inches(6.1), Inches(0.6),
             RED, corner=0.08)
    add_textbox(s, Inches(0.5), Inches(header_y + 0.13), Inches(6.1), Inches(0.4),
                "❌  THE OLD WAY", size=14, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    add_card(s, Inches(6.75), Inches(header_y), Inches(6.1), Inches(0.6),
             GREEN, corner=0.08)
    add_textbox(s, Inches(6.75), Inches(header_y + 0.13), Inches(6.1), Inches(0.4),
                "✅  FRAMEDROPS", size=14, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    y = 4.75
    row_h = 0.42
    for bad, good in rows:
        # Bad
        add_card(s, Inches(0.5), Inches(y), Inches(6.1), Inches(row_h),
                 OFF_WHITE, border=SLATE_200, corner=0.15)
        add_textbox(s, Inches(0.75), Inches(y + 0.05), Inches(5.8), Inches(0.35),
                    bad, size=11.5, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        # Good
        add_card(s, Inches(6.75), Inches(y), Inches(6.1), Inches(row_h),
                 PURPLE_LAVENDER, border=PURPLE, corner=0.15)
        add_textbox(s, Inches(7.0), Inches(y + 0.05), Inches(5.8), Inches(0.35),
                    good, size=11.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + 0.05

    add_page_number(s, 3, 10)


def slide_4_how(prs):
    s = blank(prs)
    fill_background(s, OFF_WHITE)
    add_side_rail(s, PURPLE)
    add_brand_lockup(s)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "03", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.5),
             "HOW IT WORKS", PURPLE, WHITE, size=11)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(10), Inches(1.2),
                "Shoot to selection in 3 simple steps.",
                size=44, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.9), Inches(10), Inches(0.5),
                "No uploads to Drive. No paper. No decoding.",
                size=16, color=SLATE_500)

    steps = [
        ("01", "Upload",
         "Drag-drop your shoot.\nAuto-compressed in browser.\n25 MB → 250 KB.",
         "📤"),
        ("02", "Share",
         "One branded gallery link\nvia WhatsApp.\nClient taps to view.",
         "🔗"),
        ("03", "Get selections",
         "Client taps the heart\non favorites. You see the list —\nno paper, no decoding.",
         "❤️"),
    ]

    card_y = 4.0
    cw = 4.0
    gap = 0.3
    start_x = 0.5
    for idx, (num, title, body, emoji) in enumerate(steps):
        x = Inches(start_x + idx * (cw + gap))
        # Card
        add_card(s, x, Inches(card_y), Inches(cw), Inches(3.0),
                 WHITE, border=SLATE_200, corner=0.05)
        # Top accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, Inches(card_y), Inches(cw), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PURPLE
        bar.line.fill.background()

        # Step number — large, ghosted
        add_textbox(s, x + Inches(0.3), Inches(card_y + 0.25), Inches(cw - 0.6),
                    Inches(1.2), num, size=58, color=PURPLE_LAVENDER, bold=True,
                    line_spacing=1.0)
        # Emoji circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      x + Inches(cw - 1.0), Inches(card_y + 0.4),
                                      Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.fill.background()
        add_textbox(s, x + Inches(cw - 1.0), Inches(card_y + 0.45), Inches(0.7),
                    Inches(0.7), emoji, size=24, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title
        add_textbox(s, x + Inches(0.3), Inches(card_y + 1.5), Inches(cw - 0.6),
                    Inches(0.6), title, size=24, color=NAVY, bold=True)
        # Body
        add_textbox(s, x + Inches(0.3), Inches(card_y + 2.1), Inches(cw - 0.6),
                    Inches(1.0), body, size=12, color=SLATE_500, line_spacing=1.4)

    add_page_number(s, 4, 10)


def slide_5_india(prs):
    s = blank(prs)
    fill_background(s, NAVY)
    add_side_rail(s, MARIGOLD)
    add_brand_lockup(s, on_dark=True)

    # Decorative glows
    add_glow(s, Inches(-3), Inches(2), Inches(8), INDIGO)
    add_glow(s, Inches(8), Inches(-2), Inches(6), DEEP_PURPLE)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "04", size=72, color=RGBColor(0x1f, 0x2a, 0x44), bold=True,
                align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(3.0), Inches(0.5),
             "BUILT IN INDIA  🇮🇳", MARIGOLD, NAVY, size=11)

    # Left side: huge headline
    add_textbox(s, Inches(0.5), Inches(1.8), Inches(6.5), Inches(3),
                "Made for the\nIndian wedding.",
                size=58, color=WHITE, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(6.5), Inches(2),
                "Not a US tool with INR pricing taped on.\n"
                "A product designed around how Indian\n"
                "photographers and clients actually work.",
                size=15, color=SLATE_400, line_spacing=1.5)

    # Right side: perk cards in 2-col grid
    perks = [
        ("🗣️", "Three languages", "English · Telugu · Hindi"),
        ("💳", "INR + GST native", "Razorpay UPI checkout"),
        ("📱", "Mobile-first", "Your clients live on phones"),
        ("🏦", "Wallet payouts", "Any Indian bank via NEFT"),
        ("💸", "No monthly tax", "Pay only when you publish"),
        ("🔐", "DPDP-aware", "Your data stays yours"),
    ]
    grid_x = 7.4
    grid_y = 1.8
    cw, ch = 2.7, 1.6
    gap = 0.15
    for idx, (emoji, head, sub) in enumerate(perks):
        col = idx % 2
        row = idx // 2
        x = Inches(grid_x + col * (cw + gap))
        y = Inches(grid_y + row * (ch + gap))
        add_card(s, x, y, Inches(cw), Inches(ch),
                 RGBColor(0x1a, 0x24, 0x3d), border=RGBColor(0x2c, 0x3b, 0x60),
                 corner=0.1)
        add_textbox(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                    emoji, size=18, color=WHITE)
        add_textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(cw - 0.4),
                    Inches(0.4), head, size=13, color=WHITE, bold=True)
        add_textbox(s, x + Inches(0.2), y + Inches(1.05), Inches(cw - 0.4),
                    Inches(0.5), sub, size=10, color=SLATE_400, line_spacing=1.3)

    add_page_number(s, 5, 10, on_dark=True)


def slide_6_pricing(prs):
    s = blank(prs)
    fill_background(s, WHITE)
    add_side_rail(s, PURPLE)
    add_brand_lockup(s)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "05", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    # Header strip
    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.0), Inches(0.5),
             "PRICING", PURPLE, WHITE, size=11)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(8), Inches(1.0),
                "Simple, honest pricing.",
                size=48, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.7), Inches(8), Inches(0.5),
                "Pay only when you publish. No monthly fees. No surprises.",
                size=15, color=SLATE_500)

    # Launch banner on the right
    banner = add_card(s, Inches(8.8), Inches(1.05), Inches(4.0), Inches(1.65),
                      PURPLE_LAVENDER, border=PURPLE, corner=0.06)
    add_pill(s, Inches(9.0), Inches(1.2), Inches(1.7), Inches(0.4),
             "🚀 LAUNCH OFFER", PURPLE, WHITE, size=10)
    add_textbox(s, Inches(9.0), Inches(1.7), Inches(3.6), Inches(0.6),
                "Locked for the first 100\nphotographers.",
                size=12, color=NAVY, bold=True, line_spacing=1.3)
    add_textbox(s, Inches(9.0), Inches(2.35), Inches(3.6), Inches(0.4),
                "Prices increase after.",
                size=11, color=SLATE_500)

    # Free tier strip
    add_card(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.7),
             NAVY, corner=0.08)
    add_textbox(s, Inches(0.8), Inches(3.5), Inches(8), Inches(0.5),
                "🎁  First 300 images — FREE for life",
                size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(8.5), Inches(3.5), Inches(4.0), Inches(0.5),
                "No credit card. 2-minute signup.",
                size=12, color=SLATE_400, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)

    # Three tier cards — side by side, premium look
    tiers = [
        ("STARTER", "149", "199", "1,000", "0.15", False, "Pre-wedding\nor reception"),
        ("WEDDING\nSWEET-SPOT", "229", "299", "2,000", "0.11", True, "Full wedding\nday"),
        ("WEDDING PRO", "299", "399", "3,000", "0.10", False, "Big-fat\nmulti-day"),
    ]

    card_y = 4.4
    cw = 4.0
    gap = 0.15
    start_x = 0.5
    for idx, (name, price, was, imgs, per_img, popular, scenario) in enumerate(tiers):
        x = Inches(start_x + idx * (cw + gap))
        if popular:
            add_card(s, x, Inches(card_y - 0.15), Inches(cw), Inches(3.0),
                     PURPLE_LAVENDER, border=PURPLE, border_width=2.5, corner=0.05)
            # Top badge tab
            tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x, Inches(card_y - 0.15),
                                       Inches(cw), Inches(0.35))
            tab.fill.solid()
            tab.fill.fore_color.rgb = PURPLE
            tab.line.fill.background()
            add_textbox(s, x, Inches(card_y - 0.1), Inches(cw), Inches(0.3),
                        "★  MOST POPULAR", size=10, color=WHITE, bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            content_top = card_y + 0.3
        else:
            add_card(s, x, Inches(card_y), Inches(cw), Inches(2.85),
                     OFF_WHITE, border=SLATE_200, corner=0.05)
            content_top = card_y + 0.2

        # Tier name
        add_textbox(s, x + Inches(0.3), Inches(content_top), Inches(cw - 0.6),
                    Inches(0.6), name, size=11,
                    color=PURPLE if popular else SLATE_500, bold=True,
                    line_spacing=1.1)
        # Price + strike
        add_textbox(s, x + Inches(0.3), Inches(content_top + 0.7),
                    Inches(2.0), Inches(0.9), f"₹{price}", size=46,
                    color=PURPLE if popular else NAVY, bold=True,
                    line_spacing=1.0)
        add_strike_price(s, x + Inches(2.4), Inches(content_top + 1.0),
                          Inches(1.4), f"₹{was}", size=18, color=SLATE_400)
        # Range
        add_textbox(s, x + Inches(0.3), Inches(content_top + 1.55),
                    Inches(cw - 0.6), Inches(0.4),
                    f"Up to {imgs} photos", size=13, color=NAVY, bold=True)
        # Per-image
        add_textbox(s, x + Inches(0.3), Inches(content_top + 1.9),
                    Inches(cw - 0.6), Inches(0.4),
                    f"₹{per_img} per image",
                    size=11, color=PURPLE, bold=True)
        # Scenario footer
        add_textbox(s, x + Inches(0.3), Inches(content_top + 2.3),
                    Inches(cw - 0.6), Inches(0.5),
                    scenario, size=10, color=SLATE_500, line_spacing=1.3)

    add_page_number(s, 6, 10)


def slide_7_premium(prs):
    s = blank(prs)
    fill_background(s, NAVY)
    add_side_rail(s, MARIGOLD)
    add_brand_lockup(s, on_dark=True)

    add_glow(s, Inches(8), Inches(-2), Inches(9), DEEP_PURPLE)
    add_glow(s, Inches(-2), Inches(5), Inches(7), INDIGO)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "06", size=72, color=RGBColor(0x1f, 0x2a, 0x44), bold=True,
                align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.5),
             "LOOK PREMIUM", MARIGOLD, NAVY, size=11)

    add_textbox(s, Inches(0.5), Inches(1.8), Inches(7), Inches(3.5),
                "Look like a\n₹2-lakh\nphotographer.",
                size=64, color=WHITE, bold=True, line_spacing=1.0)

    add_textbox(s, Inches(0.5), Inches(5.2), Inches(7), Inches(2),
                "Even on your very first paid shoot.\n"
                "Even when you're learning the business side.",
                size=15, color=SLATE_400, line_spacing=1.4)

    # Right: feature checklist with check icons
    perks = [
        "Custom studio branding on every gallery",
        "Watermarked, screenshot-protected previews",
        "Anti-download UI keeps your originals safe",
        "Clients see YOUR name — not 'Framedrops'",
        "Looks pro on day one",
    ]
    right_x = Inches(8.0)
    add_card(s, right_x, Inches(1.8), Inches(4.8), Inches(5.4),
             RGBColor(0x1a, 0x24, 0x3d), border=RGBColor(0x2c, 0x3b, 0x60),
             corner=0.04)
    add_textbox(s, right_x + Inches(0.4), Inches(2.0), Inches(4.0), Inches(0.5),
                "INCLUDED ON EVERY ALBUM", size=10, color=MARIGOLD, bold=True)

    y = 2.6
    for line in perks:
        chk = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  right_x + Inches(0.4), Inches(y + 0.05),
                                  Inches(0.3), Inches(0.3))
        chk.fill.solid()
        chk.fill.fore_color.rgb = PURPLE
        chk.line.fill.background()
        add_textbox(s, right_x + Inches(0.4), Inches(y + 0.0),
                    Inches(0.3), Inches(0.35), "✓", size=12, color=WHITE,
                    bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, right_x + Inches(0.85), Inches(y + 0.02),
                    Inches(3.8), Inches(0.5), line, size=13, color=WHITE,
                    line_spacing=1.3)
        y += 0.85

    add_page_number(s, 7, 10, on_dark=True)


def slide_8_value(prs):
    s = blank(prs)
    fill_background(s, OFF_WHITE)
    add_side_rail(s, GREEN)
    add_brand_lockup(s)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "07", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.5),
             "THE NUMBERS", GREEN, WHITE, size=11)

    add_textbox(s, Inches(0.5), Inches(1.7), Inches(11), Inches(1.2),
                "₹149 pays for itself in 15 minutes.",
                size=44, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.9), Inches(11), Inches(0.5),
                "Real numbers from photographers using gallery-selection tools.",
                size=15, color=SLATE_500)

    stats = [
        ("6+", "hours", "Saved per wedding —\nno more decoding lists"),
        ("0", "minutes", "Spent searching\nDrive folder-by-folder"),
        ("48", "hours", "From shoot to selection\n(vs 7+ days)"),
        ("₹500–₹2k", "per hour", "What your time is\nactually worth"),
    ]
    card_y = 4.0
    cw = 3.0
    gap = 0.13
    start_x = 0.5
    for idx, (big, unit, sub) in enumerate(stats):
        x = Inches(start_x + idx * (cw + gap))
        add_card(s, x, Inches(card_y), Inches(cw), Inches(3.0),
                 WHITE, border=PURPLE if idx == 0 else SLATE_200,
                 border_width=2 if idx == 0 else 1, corner=0.05)
        # Top accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, Inches(card_y), Inches(cw), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PURPLE if idx == 0 else SLATE_300
        bar.line.fill.background()

        # Big number
        add_textbox(s, x + Inches(0.3), Inches(card_y + 0.5), Inches(cw - 0.6),
                    Inches(1.6), big, size=64, color=PURPLE if idx == 0 else NAVY,
                    bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Unit
        add_textbox(s, x + Inches(0.3), Inches(card_y + 1.9), Inches(cw - 0.6),
                    Inches(0.4), unit, size=12, color=SLATE_500, bold=True,
                    align=PP_ALIGN.CENTER)
        # Divider
        div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(cw * 0.3), Inches(card_y + 2.3),
                                   Inches(cw * 0.4), Inches(0.02))
        div.fill.solid()
        div.fill.fore_color.rgb = SLATE_200
        div.line.fill.background()
        # Sub
        add_textbox(s, x + Inches(0.3), Inches(card_y + 2.4), Inches(cw - 0.6),
                    Inches(0.6), sub, size=10.5, color=SLATE_500,
                    align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_page_number(s, 8, 10)


def slide_9_vs(prs):
    s = blank(prs)
    fill_background(s, WHITE)
    add_side_rail(s, PURPLE)
    add_brand_lockup(s)

    add_textbox(s, Inches(11.7), Inches(0.4), Inches(1.5), Inches(2),
                "08", size=72, color=SLATE_200, bold=True, align=PP_ALIGN.RIGHT)

    add_pill(s, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.5),
             "COMPARISON", PURPLE, WHITE, size=11)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(8), Inches(1.2),
                "Why not just free cloud + paper?",
                size=40, color=NAVY, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(0.5), Inches(2.85), Inches(8), Inches(0.5),
                "Because 'free' has a hidden cost.",
                size=15, color=SLATE_500)

    # Two-column big cards
    card_y = 4.0
    card_h = 3.2

    # Left card — The Old Way
    add_card(s, Inches(0.5), Inches(card_y), Inches(6.1), Inches(card_h),
             OFF_WHITE, border=SLATE_200, corner=0.04)
    add_textbox(s, Inches(0.8), Inches(card_y + 0.25), Inches(5.5), Inches(0.4),
                "THE OLD WAY", size=11, color=SLATE_500, bold=True)
    add_textbox(s, Inches(0.8), Inches(card_y + 0.65), Inches(5.5), Inches(0.7),
                "Free cloud storage\n+ paper lists",
                size=24, color=NAVY, bold=True, line_spacing=1.1)
    add_textbox(s, Inches(0.8), Inches(card_y + 1.85), Inches(5.5), Inches(1.3),
                "❌  Clients write file names on paper\n"
                "❌  You decode handwriting, file-by-file\n"
                "❌  No branding — looks like a hobbyist\n"
                "❌  Cloud storage nags for paid upgrade\n"
                "❌  Misread one name → angry client",
                size=12, color=SLATE_500, line_spacing=1.55)

    # Right card — Framedrops
    add_card(s, Inches(6.75), Inches(card_y - 0.1), Inches(6.1), Inches(card_h + 0.2),
             PURPLE_LAVENDER, border=PURPLE, border_width=2.5, corner=0.04)
    add_pill(s, Inches(7.05), Inches(card_y), Inches(2.0), Inches(0.4),
             "★  RECOMMENDED", PURPLE, WHITE, size=10)
    add_textbox(s, Inches(7.05), Inches(card_y + 0.55), Inches(5.5), Inches(0.4),
                "FRAMEDROPS", size=11, color=PURPLE, bold=True)
    add_textbox(s, Inches(7.05), Inches(card_y + 0.95), Inches(5.5), Inches(0.7),
                "₹149 per album",
                size=32, color=PURPLE, bold=True, line_spacing=1.0)
    add_textbox(s, Inches(7.05), Inches(card_y + 1.55), Inches(5.5), Inches(0.4),
                "(only when you publish)", size=11, color=SLATE_500)

    add_textbox(s, Inches(7.05), Inches(card_y + 1.95), Inches(5.5), Inches(1.5),
                "✅  Pay zero when you don't shoot\n"
                "✅  INR-native — no FX surprises\n"
                "✅  English, Telugu, Hindi\n"
                "✅  GST invoices + Razorpay UPI\n"
                "✅  Built in India, for Indian weddings",
                size=12, color=NAVY, line_spacing=1.55)

    add_page_number(s, 9, 10)


def slide_10_cta(prs):
    s = blank(prs)
    fill_background(s, NAVY)
    add_side_rail(s, MARIGOLD)
    add_brand_lockup(s, on_dark=True)

    # Big atmospheric glows
    add_glow(s, Inches(7), Inches(-3), Inches(10), DEEP_PURPLE)
    add_glow(s, Inches(-3), Inches(4), Inches(8), INDIGO)
    add_glow(s, Inches(10), Inches(4), Inches(5), PURPLE)

    # Diagonal marigold accent
    stripe = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                  Inches(0), Inches(0.7), Inches(2.5), Inches(0.3))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = MARIGOLD
    stripe.line.fill.background()

    add_pill(s, Inches(0.5), Inches(1.5), Inches(2.8), Inches(0.55),
             "✨ START FREE TODAY", PURPLE, WHITE, size=12)

    add_textbox(s, Inches(0.5), Inches(2.3), Inches(8.5), Inches(3.5),
                "Your next wedding\ndeserves better.",
                size=64, color=WHITE, bold=True, line_spacing=1.0)

    # Trust line
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(8), Inches(0.6),
                "🎁  300 free images  ·  💳  No credit card  ·  ⚡  2-min signup",
                size=16, color=SLATE_400)

    # Tagline
    add_textbox(s, Inches(0.5), Inches(5.8), Inches(7), Inches(0.6),
                "Stop chasing. Start delivering.",
                size=18, color=MARIGOLD, bold=True)

    # Big CTA URL card — right side, prominent
    cta_x = Inches(9.3)
    cta_y = Inches(2.8)
    # Back shadow card
    add_card(s, cta_x + Inches(0.15), cta_y + Inches(0.15), Inches(3.5),
             Inches(2.0), PURPLE, corner=0.06)
    # Front card
    add_card(s, cta_x, cta_y, Inches(3.5), Inches(2.0), WHITE, corner=0.06)
    add_textbox(s, cta_x, cta_y + Inches(0.25), Inches(3.5), Inches(0.4),
                "TRY FRAMEDROPS FREE", size=11, color=PURPLE, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, cta_x, cta_y + Inches(0.7), Inches(3.5), Inches(0.8),
                WEBSITE, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # Underline accent
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cta_x + Inches(1.0), cta_y + Inches(1.4),
                                     Inches(1.5), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = PURPLE
    underline.line.fill.background()
    add_textbox(s, cta_x, cta_y + Inches(1.55), Inches(3.5), Inches(0.4),
                "Open on phone or laptop", size=10, color=SLATE_500,
                align=PP_ALIGN.CENTER)

    add_page_number(s, 10, 10, on_dark=True)


# ────────────────────────────────────────────────────────────────────────────


def build() -> None:
    prs = make_presentation()
    slide_1_hook(prs)
    slide_2_pain(prs)
    slide_3_compare(prs)
    slide_4_how(prs)
    slide_5_india(prs)
    slide_6_pricing(prs)
    slide_7_premium(prs)
    slide_8_value(prs)
    slide_9_vs(prs)
    slide_10_cta(prs)
    prs.save(OUT_FILE)
    print(f"✓ Wrote {OUT_FILE}")


if __name__ == "__main__":
    build()
