"""
Framedrops — Portrait (9:16) Premium Pitch Deck for Instagram Reels / YouTube Shorts
Generates a PowerPoint file with modern SaaS/gradient styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------- Palette ----------
NAVY        = RGBColor(0x08, 0x1B, 0x3A)
NAVY_DEEP   = RGBColor(0x05, 0x10, 0x26)
ROYAL       = RGBColor(0x00, 0x57, 0xFF)
PURPLE      = RGBColor(0x7B, 0x2C, 0xFF)
PINK        = RGBColor(0xFF, 0x4D, 0x8D)
ORANGE      = RGBColor(0xFF, 0x8A, 0x3D)
SOFT_WHITE  = RGBColor(0xF8, 0xFA, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0xB6, 0xC2, 0xE0)
CARD_FILL   = RGBColor(0x10, 0x24, 0x4A)
DIVIDER     = RGBColor(0x1B, 0x2E, 0x5C)

# ---------- Setup: 9:16 Portrait ----------
prs = Presentation()
# Instagram Reels / Shorts: 1080x1920 -> use 7.5 x 13.333 inches (same aspect)
prs.slide_width  = Inches(7.5)
prs.slide_height = Inches(13.333)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill_rgb, line=False, line_rgb=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb or WHITE
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_gradient_rect(slide, x, y, w, h, stops, angle=2700000, shape=MSO_SHAPE.RECTANGLE, line=False, line_rgb=None, line_w=0.75):
    """
    stops: list of (position 0..1, RGBColor)
    angle in 60000ths of a degree (PowerPoint convention). 2700000 = 45 deg.
    Common: 0 = left->right, 5400000 = top->bottom, 2700000 = TL->BR (45deg), 16200000 = BL->TR.
    """
    s = slide.shapes.add_shape(shape, x, y, w, h)
    sp = s.fill._xPr  # spPr
    # remove existing fill
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)

    grad = etree.SubElement(sp, qn('a:gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '0')
    tileRect = etree.SubElement(grad, qn('a:tileRect'))

    # line
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb or WHITE
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_gradient_line_rect(slide, x, y, w, h, stops, line_w_pt=1.5, angle=0):
    """Rounded rect with NO fill and a gradient stroke (neon outline)."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.background()
    sp = s.line._get_or_add_ln() if False else s._element.spPr  # use spPr
    # Remove existing line
    for tag in ('a:ln',):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    ln = etree.SubElement(sp, qn('a:ln'))
    ln.set('w', str(int(line_w_pt * 12700)))
    ln.set('cap', 'rnd')
    grad = etree.SubElement(ln, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '0')
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font='Helvetica Neue', italic=False, line_spacing=1.15, letter_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if letter_sp is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(letter_sp))
    return tb


def add_gradient_text(slide, x, y, w, h, text, size=44, bold=True, stops=None, align=PP_ALIGN.LEFT,
                      anchor=MSO_ANCHOR.TOP, font='Helvetica Neue', angle=0, line_spacing=1.05, letter_sp=None):
    """Text fill with linear gradient."""
    if stops is None:
        stops = [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold

        rPr = r._r.get_or_add_rPr()
        # remove existing fills on run
        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        grad = etree.SubElement(rPr, qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        for pos, rgb in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(int(pos * 100000)))
            srgb = etree.SubElement(gs, qn('a:srgbClr'))
            srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(angle))
        lin.set('scaled', '0')

        if letter_sp is not None:
            rPr.set('spc', str(letter_sp))
    return tb


def add_blob(slide, cx, cy, r, rgb, alpha=40000, shape=MSO_SHAPE.OVAL):
    """Soft glowing blob (translucent oval)."""
    x = cx - r // 2
    y = cy - r // 2
    s = slide.shapes.add_shape(shape, x, y, r, r)
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    sf = etree.SubElement(sp, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    alphaEl = etree.SubElement(srgb, qn('a:alpha'))
    alphaEl.set('val', str(alpha))
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def background(slide, variant='dark'):
    """Paint dark navy background + atmospheric blobs."""
    if variant == 'dark':
        add_rect(slide, 0, 0, SW, SH, NAVY_DEEP)
        # subtle vertical gradient overlay
        add_gradient_rect(slide, 0, 0, SW, SH,
                          stops=[(0.0, NAVY_DEEP), (1.0, NAVY)],
                          angle=5400000)
        # glow blobs
        add_blob(slide, Inches(1.0), Inches(1.2), Inches(4.5), PURPLE, alpha=30000)
        add_blob(slide, Inches(6.8), Inches(3.5), Inches(4.0), ROYAL, alpha=28000)
        add_blob(slide, Inches(0.6), Inches(10.5), Inches(4.8), PINK, alpha=22000)
        add_blob(slide, Inches(7.0), Inches(12.0), Inches(3.8), ORANGE, alpha=18000)
    elif variant == 'light':
        add_rect(slide, 0, 0, SW, SH, SOFT_WHITE)
        add_blob(slide, Inches(0.2), Inches(0.6), Inches(5.0), ROYAL, alpha=20000)
        add_blob(slide, Inches(7.2), Inches(4.0), Inches(4.6), PURPLE, alpha=18000)
        add_blob(slide, Inches(1.0), Inches(11.0), Inches(5.2), PINK, alpha=18000)
    elif variant == 'split':
        # top dark navy, bottom soft white
        add_rect(slide, 0, 0, SW, Inches(8.2), NAVY_DEEP)
        add_gradient_rect(slide, 0, 0, SW, Inches(8.2),
                          stops=[(0.0, NAVY_DEEP), (1.0, NAVY)], angle=5400000)
        add_rect(slide, 0, Inches(8.2), SW, SH - Inches(8.2), SOFT_WHITE)
        add_blob(slide, Inches(6.5), Inches(2.0), Inches(4.2), PURPLE, alpha=32000)
        add_blob(slide, Inches(1.0), Inches(6.5), Inches(4.0), PINK, alpha=24000)


def gradient_divider(slide, y, w_in=2.6, x_center=None, thickness_pt=3):
    if x_center is None:
        x_center = SW // 2
    w = Inches(w_in)
    x = x_center - w // 2
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.06))
    # gradient fill
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', '0'); lin.set('scaled', '0')
    s.line.fill.background()
    s.shadow.inherit = False


def pill(slide, x, y, w, h, label, fill=None, gradient_stops=None, text_color=WHITE, size=11, bold=True, letter_sp=200):
    if gradient_stops:
        add_gradient_rect(slide, x, y, w, h, gradient_stops, angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    else:
        add_rect(slide, x, y, w, h, fill or CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = add_text(slide, x, y, w, h, label, size=size, bold=bold, color=text_color,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=letter_sp)
    return tb


def glass_card(slide, x, y, w, h, fill_alpha=55000, line_alpha=50000, corner=None):
    """Glassmorphism panel: translucent fill + soft white outline."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # translucent white fill
    sp = s.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    sf = etree.SubElement(sp, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', 'FFFFFF')
    alphaEl = etree.SubElement(srgb, qn('a:alpha'))
    alphaEl.set('val', str(fill_alpha))
    # outline
    s.line.color.rgb = WHITE
    s.line.width = Pt(0.75)
    lnEl = s.line._get_or_add_ln()
    # set alpha on line color
    solid = lnEl.find(qn('a:solidFill'))
    if solid is not None:
        srgb2 = solid.find(qn('a:srgbClr'))
        if srgb2 is not None:
            for child in list(srgb2):
                srgb2.remove(child)
            a = etree.SubElement(srgb2, qn('a:alpha'))
            a.set('val', str(line_alpha))
    s.shadow.inherit = False
    return s


def neon_card(slide, x, y, w, h, body_fill=CARD_FILL, stops=None):
    """Filled rounded card + gradient outline overlay."""
    if stops is None:
        stops = [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]
    add_rect(slide, x, y, w, h, body_fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(slide, x, y, w, h, stops, line_w_pt=1.75, angle=2700000)


def wordmark(slide, x, y, size=18):
    """Frame● drops wordmark (white text + gradient dot)."""
    tb = slide.shapes.add_textbox(x, y, Inches(2.2), Inches(0.45))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = 'Frame'
    r1.font.name = 'Helvetica Neue'; r1.font.size = Pt(size); r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = 'drops'
    r2.font.name = 'Helvetica Neue'; r2.font.size = Pt(size); r2.font.bold = True
    r2.font.color.rgb = WHITE
    # gradient dot after
    add_gradient_rect(slide, x + Inches(1.85), y + Inches(0.13), Inches(0.18), Inches(0.18),
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=2700000, shape=MSO_SHAPE.OVAL)


def slide_chrome(slide, idx, total):
    """Top wordmark + slide counter, consistent on every slide."""
    wordmark(slide, Inches(0.45), Inches(0.5), size=16)
    # counter pill
    pw, ph = Inches(0.95), Inches(0.36)
    px = SW - pw - Inches(0.45)
    py = Inches(0.54)
    add_rect(slide, px, py, pw, ph, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, px, py, pw, ph, f'{idx:02d} / {total:02d}',
             size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=300)


def bottom_brand(slide, color=MUTED):
    add_text(slide, 0, SH - Inches(0.55), SW, Inches(0.3),
             'framedrops.in',
             size=11, bold=True, color=color, align=PP_ALIGN.CENTER, letter_sp=400)


# ---------- ICONS (built from primitive shapes) ----------
def icon_box(slide, x, y, size_in=0.95, gradient=True, stops=None):
    """Rounded square icon container with gradient or filled background."""
    if stops is None:
        stops = [(0.0, ROYAL), (1.0, PURPLE)]
    s_in = size_in
    if gradient:
        add_gradient_rect(slide, x, y, Inches(s_in), Inches(s_in), stops,
                          angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    else:
        add_rect(slide, x, y, Inches(s_in), Inches(s_in), CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    return (x, y, Inches(s_in))


def icon_camera(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops)
    # body
    pad = Inches(size_in * 0.18)
    inner_x = x + pad
    inner_y = y + Inches(size_in * 0.32)
    inner_w = Inches(size_in) - 2 * pad
    inner_h = Inches(size_in * 0.42)
    add_rect(slide, inner_x, inner_y, inner_w, inner_h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # lens
    lens_d = Inches(size_in * 0.28)
    lens_x = x + Inches(size_in) / 2 - lens_d / 2
    lens_y = inner_y + inner_h / 2 - lens_d / 2
    add_rect(slide, lens_x, lens_y, lens_d, lens_d, NAVY_DEEP, shape=MSO_SHAPE.OVAL)
    add_rect(slide, lens_x + lens_d * 0.25, lens_y + lens_d * 0.25, lens_d * 0.5, lens_d * 0.5, WHITE, shape=MSO_SHAPE.OVAL)
    # top viewfinder bump
    vw = Inches(size_in * 0.22); vh = Inches(size_in * 0.08)
    vx = x + Inches(size_in) * 0.28; vy = inner_y - vh
    add_rect(slide, vx, vy, vw, vh, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def icon_heart(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, PINK), (1.0, ORANGE)])
    s = slide.shapes.add_shape(MSO_SHAPE.HEART,
                               x + Inches(size_in * 0.18),
                               y + Inches(size_in * 0.22),
                               Inches(size_in * 0.64),
                               Inches(size_in * 0.56))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background(); s.shadow.inherit = False


def icon_cloud(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops)
    s = slide.shapes.add_shape(MSO_SHAPE.CLOUD,
                               x + Inches(size_in * 0.12),
                               y + Inches(size_in * 0.28),
                               Inches(size_in * 0.76),
                               Inches(size_in * 0.5))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background(); s.shadow.inherit = False
    # up arrow
    a = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                               x + Inches(size_in * 0.42),
                               y + Inches(size_in * 0.42),
                               Inches(size_in * 0.16),
                               Inches(size_in * 0.26))
    a.fill.solid(); a.fill.fore_color.rgb = ROYAL
    a.line.fill.background(); a.shadow.inherit = False


def icon_folder(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, ORANGE), (1.0, PINK)])
    # folder tab
    tab_w = Inches(size_in * 0.32); tab_h = Inches(size_in * 0.10)
    add_rect(slide, x + Inches(size_in * 0.18), y + Inches(size_in * 0.28),
             tab_w, tab_h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # folder body
    add_rect(slide, x + Inches(size_in * 0.16), y + Inches(size_in * 0.36),
             Inches(size_in * 0.68), Inches(size_in * 0.40),
             WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def icon_link(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, PURPLE), (1.0, PINK)])
    # two overlapping rounded rects to suggest a link
    for i, dx in enumerate([0.20, 0.42]):
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(size_in * dx),
                                   y + Inches(size_in * 0.36),
                                   Inches(size_in * 0.38),
                                   Inches(size_in * 0.22))
        r.fill.background(); r.line.color.rgb = WHITE; r.line.width = Pt(2.25)
        r.shadow.inherit = False
    # connector
    add_rect(slide, x + Inches(size_in * 0.43), y + Inches(size_in * 0.45),
             Inches(size_in * 0.14), Inches(size_in * 0.05), WHITE,
             shape=MSO_SHAPE.RECTANGLE)


def icon_payment(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, ROYAL), (1.0, PURPLE)])
    # card body
    add_rect(slide, x + Inches(size_in * 0.14), y + Inches(size_in * 0.28),
             Inches(size_in * 0.72), Inches(size_in * 0.44),
             WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # magnetic stripe
    add_rect(slide, x + Inches(size_in * 0.14), y + Inches(size_in * 0.38),
             Inches(size_in * 0.72), Inches(size_in * 0.07),
             NAVY_DEEP)
    # chip
    add_rect(slide, x + Inches(size_in * 0.20), y + Inches(size_in * 0.52),
             Inches(size_in * 0.16), Inches(size_in * 0.12),
             ORANGE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def icon_check(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, ROYAL), (1.0, PURPLE)])
    # big white check mark using a rotated rounded rect "tick"
    p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x + Inches(size_in * 0.22),
                                y + Inches(size_in * 0.48),
                                Inches(size_in * 0.28),
                                Inches(size_in * 0.10))
    p1.rotation = -45
    p1.fill.solid(); p1.fill.fore_color.rgb = WHITE
    p1.line.fill.background(); p1.shadow.inherit = False
    p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x + Inches(size_in * 0.42),
                                y + Inches(size_in * 0.36),
                                Inches(size_in * 0.42),
                                Inches(size_in * 0.10))
    p2.rotation = 45
    p2.fill.solid(); p2.fill.fore_color.rgb = WHITE
    p2.line.fill.background(); p2.shadow.inherit = False


def icon_album(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, PURPLE), (1.0, ROYAL)])
    # stacked photo cards
    for i, off in enumerate([(0.18, 0.30), (0.26, 0.36), (0.34, 0.42)]):
        ox, oy = off
        add_rect(slide,
                 x + Inches(size_in * ox),
                 y + Inches(size_in * oy),
                 Inches(size_in * 0.42),
                 Inches(size_in * 0.30),
                 WHITE if i == 2 else SOFT_WHITE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def icon_chat(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, PINK), (1.0, PURPLE)])
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                               x + Inches(size_in * 0.14),
                               y + Inches(size_in * 0.22),
                               Inches(size_in * 0.72),
                               Inches(size_in * 0.50))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background(); s.shadow.inherit = False
    # dots
    for i, dx in enumerate([0.30, 0.46, 0.62]):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x + Inches(size_in * dx),
                                   y + Inches(size_in * 0.42),
                                   Inches(size_in * 0.08),
                                   Inches(size_in * 0.08))
        d.fill.solid(); d.fill.fore_color.rgb = PURPLE
        d.line.fill.background(); d.shadow.inherit = False


def icon_lightning(slide, x, y, size_in=0.95, stops=None):
    icon_box(slide, x, y, size_in, stops=stops or [(0.0, ORANGE), (1.0, PINK)])
    s = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT,
                               x + Inches(size_in * 0.30),
                               y + Inches(size_in * 0.18),
                               Inches(size_in * 0.40),
                               Inches(size_in * 0.64))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background(); s.shadow.inherit = False


# ---------- Phone mockup (used as a visual element) ----------
def phone_mockup(slide, x, y, w, h, screen_content=None):
    """A dark rounded 'phone' frame with a glass screen."""
    # outer body
    add_rect(slide, x, y, w, h, NAVY_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # subtle gradient border (we fake with 2 rects)
    add_gradient_line_rect(slide, x, y, w, h, stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=2.0, angle=2700000)
    # screen
    inset = Inches(0.12)
    sx, sy = x + inset, y + inset
    sw, sh = w - 2 * inset, h - 2 * inset
    add_gradient_rect(slide, sx, sy, sw, sh,
                      stops=[(0.0, RGBColor(0x10, 0x1F, 0x42)),
                             (1.0, RGBColor(0x1A, 0x12, 0x44))],
                      angle=5400000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # notch
    nw = w * 0.32
    add_rect(slide, x + (w - nw) / 2, y + Inches(0.08), nw, Inches(0.18), NAVY_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if screen_content:
        screen_content(slide, sx, sy, sw, sh)


def screen_album_grid(slide, sx, sy, sw, sh):
    # header
    add_text(slide, sx + Inches(0.18), sy + Inches(0.34), sw - Inches(0.36), Inches(0.25),
             'Sneha & Rohan • Sangeet',
             size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, sx + Inches(0.18), sy + Inches(0.55), sw - Inches(0.36), Inches(0.2),
             '482 photos  •  Select your favorites',
             size=7, color=MUTED, align=PP_ALIGN.LEFT)
    # grid
    cols = 3
    rows = 5
    gap = Inches(0.06)
    pad = Inches(0.18)
    grid_top = sy + Inches(0.95)
    grid_w = sw - 2 * pad
    cell_w = (grid_w - (cols - 1) * gap) / cols
    cell_h = cell_w  # square
    grad_palettes = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
        [(0.0, ROYAL), (1.0, PINK)],
        [(0.0, ORANGE), (1.0, PURPLE)],
        [(0.0, ROYAL), (1.0, ORANGE)],
    ]
    hearted = {(0, 1), (1, 0), (2, 2), (3, 1), (4, 0)}
    for r in range(rows):
        for c in range(cols):
            cx = sx + pad + c * (cell_w + gap)
            cy = grid_top + r * (cell_h + gap)
            pal = grad_palettes[(r * cols + c) % len(grad_palettes)]
            add_gradient_rect(slide, cx, cy, cell_w, cell_h, pal, angle=2700000,
                              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            if (r, c) in hearted:
                # heart badge bottom-right
                hsize = Inches(0.22)
                hs = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            cx + cell_w - hsize - Inches(0.04),
                                            cy + cell_h - hsize - Inches(0.04),
                                            hsize, hsize)
                hs.fill.solid(); hs.fill.fore_color.rgb = WHITE
                hs.line.fill.background(); hs.shadow.inherit = False
                hh = slide.shapes.add_shape(MSO_SHAPE.HEART,
                                            cx + cell_w - hsize - Inches(0.01),
                                            cy + cell_h - hsize - Inches(0.01),
                                            hsize * 0.78, hsize * 0.7)
                hh.fill.solid(); hh.fill.fore_color.rgb = PINK
                hh.line.fill.background(); hh.shadow.inherit = False


def screen_upload(slide, sx, sy, sw, sh):
    add_text(slide, sx + Inches(0.18), sy + Inches(0.34), sw - Inches(0.36), Inches(0.25),
             'Upload Shoot',
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, sx + Inches(0.18), sy + Inches(0.55), sw - Inches(0.36), Inches(0.2),
             'Smart compression  •  Quality preserved',
             size=7, color=MUTED, align=PP_ALIGN.LEFT)
    # upload area
    pad = Inches(0.22)
    bx = sx + pad; by = sy + Inches(0.95)
    bw = sw - 2 * pad; bh = Inches(1.4)
    add_gradient_line_rect(slide, bx, by, bw, bh,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=1.5, angle=2700000)
    icon_cloud(slide, bx + bw / 2 - Inches(0.32), by + Inches(0.25), size_in=0.65,
               stops=[(0.0, ROYAL), (1.0, PURPLE)])
    add_text(slide, bx, by + Inches(0.95), bw, Inches(0.3),
             'Drop your shoot here',
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # progress card
    py = by + bh + Inches(0.2)
    add_rect(slide, bx, py, bw, Inches(0.7), CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, bx + Inches(0.16), py + Inches(0.08), bw, Inches(0.2),
             'sangeet-RAW.zip',
             size=8, bold=True, color=WHITE)
    add_text(slide, bx + Inches(0.16), py + Inches(0.27), bw, Inches(0.2),
             '1,248 photos  •  4.7 GB → 1.9 GB',
             size=7, color=MUTED)
    # progress bar
    pbx = bx + Inches(0.16); pby = py + Inches(0.50)
    pbw = bw - Inches(0.32); pbh = Inches(0.10)
    add_rect(slide, pbx, pby, pbw, pbh, RGBColor(0x1B, 0x2E, 0x5C), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_rect(slide, pbx, pby, pbw * 0.72, pbh,
                      stops=[(0.0, ROYAL), (1.0, PINK)], angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # stats
    sy2 = py + Inches(0.95)
    add_text(slide, bx, sy2, bw, Inches(0.25),
             '⚡ Upload accelerated • Originals safe',
             size=7, color=MUTED, align=PP_ALIGN.LEFT)


def screen_link_share(slide, sx, sy, sw, sh):
    add_text(slide, sx + Inches(0.18), sy + Inches(0.34), sw - Inches(0.36), Inches(0.25),
             'Share Album',
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    pad = Inches(0.22)
    bx = sx + pad; by = sy + Inches(0.9)
    bw = sw - 2 * pad
    # link card
    add_rect(slide, bx, by, bw, Inches(0.6), CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, bx + Inches(0.16), by + Inches(0.10), bw, Inches(0.2),
             'framedrops.in/sneha-rohan',
             size=9, bold=True, color=WHITE)
    add_text(slide, bx + Inches(0.16), by + Inches(0.32), bw, Inches(0.2),
             'Private link  •  Any device',
             size=7, color=MUTED)
    # device row
    dy = by + Inches(0.85)
    devices = ['Phone', 'Tablet', 'Laptop']
    grad_set = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
    ]
    cw = (bw - Inches(0.2)) / 3
    for i, name in enumerate(devices):
        cx = bx + i * (cw + Inches(0.1))
        add_gradient_rect(slide, cx, dy, cw, Inches(1.0), grad_set[i],
                          angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, cx, dy + Inches(0.65), cw, Inches(0.25), name,
                 size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # picks
    py2 = dy + Inches(1.2)
    add_text(slide, bx, py2, bw, Inches(0.2),
             '124 favorites • 8 clients viewing',
             size=8, color=MUTED, align=PP_ALIGN.CENTER)


def screen_payment(slide, sx, sy, sw, sh):
    add_text(slide, sx + Inches(0.18), sy + Inches(0.34), sw - Inches(0.36), Inches(0.25),
             'Payments',
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    pad = Inches(0.22)
    bx = sx + pad; by = sy + Inches(0.9)
    bw = sw - 2 * pad
    # invoice card
    add_gradient_rect(slide, bx, by, bw, Inches(1.6),
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, bx + Inches(0.18), by + Inches(0.16), bw, Inches(0.22),
             'Invoice • Sneha & Rohan',
             size=8, bold=True, color=WHITE)
    add_text(slide, bx + Inches(0.18), by + Inches(0.4), bw, Inches(0.4),
             '₹ 42,000',
             size=22, bold=True, color=WHITE)
    add_text(slide, bx + Inches(0.18), by + Inches(0.88), bw, Inches(0.25),
             'Sangeet + Wedding • Final delivery',
             size=8, color=WHITE)
    # paid pill
    pill(slide, bx + bw - Inches(0.95), by + Inches(0.18), Inches(0.78), Inches(0.28),
         'PAID', gradient_stops=[(0.0, WHITE), (1.0, SOFT_WHITE)],
         text_color=PURPLE, size=9)
    # rows
    rows = [('Razorpay', '₹ 24,000', 'Today'),
            ('UPI', '₹ 12,000', 'Yesterday'),
            ('Bank transfer', '₹ 6,000', 'May 18')]
    for i, (m, amt, when) in enumerate(rows):
        ry = by + Inches(1.85) + i * Inches(0.55)
        add_rect(slide, bx, ry, bw, Inches(0.48), CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, bx + Inches(0.16), ry + Inches(0.08), Inches(1.2), Inches(0.18),
                 m, size=8, bold=True, color=WHITE)
        add_text(slide, bx + Inches(0.16), ry + Inches(0.26), Inches(1.2), Inches(0.18),
                 when, size=7, color=MUTED)
        add_text(slide, bx, ry + Inches(0.14), bw - Inches(0.16), Inches(0.22),
                 amt, size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def screen_done(slide, sx, sy, sw, sh):
    # big check
    cx = sx + sw / 2 - Inches(0.7)
    cy = sy + Inches(1.0)
    add_gradient_rect(slide, cx, cy, Inches(1.4), Inches(1.4),
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=2700000, shape=MSO_SHAPE.OVAL)
    # check mark
    p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                cx + Inches(0.35), cy + Inches(0.74),
                                Inches(0.38), Inches(0.12))
    p1.rotation = -45
    p1.fill.solid(); p1.fill.fore_color.rgb = WHITE
    p1.line.fill.background(); p1.shadow.inherit = False
    p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                cx + Inches(0.58), cy + Inches(0.52),
                                Inches(0.60), Inches(0.12))
    p2.rotation = 45
    p2.fill.solid(); p2.fill.fore_color.rgb = WHITE
    p2.line.fill.background(); p2.shadow.inherit = False
    add_text(slide, sx, cy + Inches(1.6), sw, Inches(0.4),
             'Delivered', size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sx, cy + Inches(2.0), sw, Inches(0.3),
             '1,248 originals • 124 selected', size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDES
# ============================================================

TOTAL_SLIDES = 15  # we'll count as we add

# ---------- Slide 1: Hook ----------
def slide_1_hook(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    # eyebrow pill
    pill(s, Inches(2.45), Inches(2.4), Inches(2.6), Inches(0.45),
         'FOR PHOTOGRAPHERS',
         gradient_stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
         text_color=WHITE, size=11, letter_sp=400)

    # main hook
    add_text(s, Inches(0.55), Inches(3.2), Inches(6.4), Inches(1.0),
             'Still delivering',
             size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, Inches(0.55), Inches(4.2), Inches(6.4), Inches(2.2),
                      'client photos\nthe hard way?',
                      size=68, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.05)

    # pain bullets — glass card
    gx, gy = Inches(0.7), Inches(7.6)
    gw, gh = Inches(6.1), Inches(3.6)
    glass_card(s, gx, gy, gw, gh, fill_alpha=8000, line_alpha=35000)

    pains = [
        ('WhatsApp screenshots flying around',  PINK),
        ('“Send me photo #87 again, please”',    ORANGE),
        ('Endless follow-up chats',              PURPLE),
        ('USB drives. Pen drives. Lost drives.', ROYAL),
    ]
    for i, (text, col) in enumerate(pains):
        py = gy + Inches(0.35) + i * Inches(0.78)
        # bullet dot
        add_gradient_rect(s, gx + Inches(0.35), py + Inches(0.18), Inches(0.22), Inches(0.22),
                          stops=[(0.0, col), (1.0, PURPLE)], angle=2700000, shape=MSO_SHAPE.OVAL)
        add_text(s, gx + Inches(0.75), py, gw - Inches(1.0), Inches(0.55),
                 text, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # scroll hint
    add_text(s, 0, SH - Inches(1.0), SW, Inches(0.35),
             'Swipe up',
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=600)
    add_gradient_text(s, 0, SH - Inches(0.7), SW, Inches(0.4),
                      '↓',
                      size=22, bold=True,
                      stops=[(0.0, PINK), (1.0, ORANGE)],
                      align=PP_ALIGN.CENTER)


# ---------- Slide 2: Meet Framedrops ----------
def slide_2_meet(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.85), Inches(2.5), Inches(1.8), Inches(0.45),
         'INTRODUCING',
         gradient_stops=[(0.0, PURPLE), (1.0, PINK)],
         text_color=WHITE, size=11, letter_sp=500)

    add_text(s, Inches(0.55), Inches(3.4), Inches(6.4), Inches(0.9),
             'Meet',
             size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_gradient_text(s, Inches(0.4), Inches(4.3), Inches(6.7), Inches(1.8),
                      'Framedrops.',
                      size=110, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.0,
                      letter_sp=-50)

    gradient_divider(s, Inches(6.6), w_in=2.4)

    add_text(s, Inches(0.6), Inches(7.0), Inches(6.3), Inches(1.0),
             'The all-in-one platform for\nphotographers to share, select & deliver.',
             size=18, bold=False, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.35)

    # feature chips row
    chips = [('Upload',  [(0.0, ROYAL),  (1.0, PURPLE)]),
             ('Select',  [(0.0, PURPLE), (1.0, PINK)]),
             ('Transfer',[(0.0, PINK),   (1.0, ORANGE)]),
             ('Pay',     [(0.0, ORANGE), (1.0, ROYAL)])]
    cy = Inches(8.5)
    cw, ch = Inches(1.45), Inches(0.55)
    gap = Inches(0.10)
    total_w = 4 * cw + 3 * gap
    start_x = (SW - total_w) // 2
    for i, (label, grad) in enumerate(chips):
        cx = start_x + i * (cw + gap)
        add_gradient_line_rect(s, cx, cy, cw, ch, grad, line_w_pt=1.5, angle=2700000)
        add_text(s, cx, cy, cw, ch, label, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=300)

    # phone mockup teaser
    phone_w = Inches(3.0); phone_h = Inches(2.8)
    phone_x = (SW - phone_w) // 2
    phone_y = Inches(9.6)
    phone_mockup(s, phone_x, phone_y, phone_w, phone_h, screen_content=None)
    # mini grid inside the teaser screen
    sx, sy = phone_x + Inches(0.12), phone_y + Inches(0.12)
    sw, sh = phone_w - Inches(0.24), phone_h - Inches(0.24)
    cols = 3; rows = 3; gap = Inches(0.05)
    pad = Inches(0.1)
    cw2 = (sw - 2 * pad - (cols - 1) * gap) / cols
    ch2 = cw2
    palettes = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
        [(0.0, ROYAL), (1.0, PINK)],
        [(0.0, ORANGE), (1.0, PURPLE)],
        [(0.0, ROYAL), (1.0, ORANGE)],
    ]
    for r in range(rows):
        for c in range(cols):
            ccx = sx + pad + c * (cw2 + gap)
            ccy = sy + pad + Inches(0.2) + r * (ch2 + gap)
            add_gradient_rect(s, ccx, ccy, cw2, ch2, palettes[(r * cols + c) % len(palettes)],
                              angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


# ---------- Slide 3: Upload ----------
def slide_3_upload(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    # step pill
    pill(s, Inches(0.55), Inches(1.55), Inches(1.6), Inches(0.42),
         'STEP 01',
         gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)],
         text_color=WHITE, size=10, letter_sp=500)
    add_text(s, Inches(2.25), Inches(1.55), Inches(4.5), Inches(0.42),
             'UPLOAD', size=11, bold=True, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE, letter_sp=600)

    # headline
    add_text(s, Inches(0.55), Inches(2.2), Inches(6.4), Inches(0.9),
             'Upload your shoot',
             size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.05)
    add_gradient_text(s, Inches(0.55), Inches(3.0), Inches(6.4), Inches(1.5),
                      'in seconds.',
                      size=64, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=2700000, line_spacing=1.0)

    add_text(s, Inches(0.55), Inches(4.5), Inches(6.4), Inches(1.0),
             'Smart compression keeps\nyour quality razor sharp.',
             size=18, color=MUTED, align=PP_ALIGN.LEFT, line_spacing=1.4)

    # phone mockup with upload screen
    phone_w = Inches(3.4); phone_h = Inches(5.6)
    phone_x = (SW - phone_w) // 2
    phone_y = Inches(6.1)
    phone_mockup(s, phone_x, phone_y, phone_w, phone_h, screen_content=screen_upload)

    # stat cards under phone
    stats = [('10×', 'Faster uploads'), ('100%', 'Quality preserved'), ('GB → MB', 'Smart compress')]
    sy_row = Inches(11.95)
    cw = Inches(2.0); ch = Inches(0.95); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    start_x = (SW - total_w) // 2
    grads = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
    ]
    for i, (val, label) in enumerate(stats):
        cx = start_x + i * (cw + gap)
        add_gradient_line_rect(s, cx, sy_row, cw, ch, grads[i], line_w_pt=1.5, angle=2700000)
        add_gradient_text(s, cx, sy_row + Inches(0.1), cw, Inches(0.45), val,
                          size=20, bold=True, stops=grads[i],
                          align=PP_ALIGN.CENTER, angle=0)
        add_text(s, cx, sy_row + Inches(0.55), cw, Inches(0.35), label,
                 size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=200)


# ---------- Slide 4: Albums ----------
def slide_4_albums(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(0.55), Inches(1.55), Inches(1.6), Inches(0.42),
         'STEP 02',
         gradient_stops=[(0.0, PURPLE), (1.0, PINK)],
         text_color=WHITE, size=10, letter_sp=500)
    add_text(s, Inches(2.25), Inches(1.55), Inches(4.5), Inches(0.42),
             'ORGANIZE', size=11, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, letter_sp=600)

    add_text(s, Inches(0.55), Inches(2.2), Inches(6.4), Inches(0.9),
             'Beautiful albums for',
             size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_gradient_text(s, Inches(0.55), Inches(2.95), Inches(6.4), Inches(2.3),
                      'every client,\nevent, season.',
                      size=52, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=2700000, line_spacing=1.05)

    # album cards grid
    albums = [
        ('Sneha & Rohan', 'Wedding · 1,248 photos', [(0.0, ROYAL), (1.0, PURPLE)], icon_camera),
        ('Aarav Turns 5', 'Birthday · 312 photos',  [(0.0, PURPLE),(1.0, PINK)],   icon_album),
        ('Verma Family',  'Maternity · 184 photos', [(0.0, PINK),  (1.0, ORANGE)], icon_heart),
        ('Diwali 2026',   'Brand · 96 photos',      [(0.0, ORANGE),(1.0, ROYAL)],  icon_lightning),
    ]
    cy = Inches(6.4)
    cw = Inches(3.15); ch = Inches(2.4)
    gap = Inches(0.2)
    grid_w = 2 * cw + gap
    sx = (SW - grid_w) // 2
    for i, (title, sub, grad, icon_fn) in enumerate(albums):
        col = i % 2
        row = i // 2
        cx = sx + col * (cw + gap)
        cyy = cy + row * (ch + gap)
        # card body
        add_rect(s, cx, cyy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, cx, cyy, cw, ch, grad, line_w_pt=1.5, angle=2700000)
        # top gradient banner inside card
        add_gradient_rect(s, cx + Inches(0.15), cyy + Inches(0.15), cw - Inches(0.3), Inches(1.1),
                          stops=grad, angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # icon
        icon_fn(s, cx + Inches(0.25), cyy + Inches(0.32), size_in=0.78, stops=grad)
        # title/sub
        add_text(s, cx + Inches(1.2), cyy + Inches(0.4), cw - Inches(1.3), Inches(0.4),
                 title, size=12, bold=True, color=WHITE)
        add_text(s, cx + Inches(1.2), cyy + Inches(0.75), cw - Inches(1.3), Inches(0.3),
                 sub, size=8, color=SOFT_WHITE)
        # mini photo strip
        strip_y = cyy + Inches(1.45)
        strip_pad = Inches(0.18)
        sw2 = cw - 2 * strip_pad
        thumb_w = (sw2 - Inches(0.16) * 3) / 4
        thumb_h = Inches(0.65)
        for j in range(4):
            tx = cx + strip_pad + j * (thumb_w + Inches(0.05))
            sub_grad = [(0.0, grad[0][1]), (1.0, grad[-1][1])]
            # rotate hue a bit
            add_gradient_rect(s, tx, strip_y, thumb_w, thumb_h, sub_grad, angle=2700000,
                              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # footer line: view
        add_text(s, cx + Inches(0.2), cyy + ch - Inches(0.35), cw - Inches(0.4), Inches(0.25),
                 'View album →', size=9, bold=True, color=MUTED, letter_sp=200)


# ---------- Slide 5: Share one link ----------
def slide_5_share(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(0.55), Inches(1.55), Inches(1.6), Inches(0.42),
         'STEP 03',
         gradient_stops=[(0.0, PINK), (1.0, ORANGE)],
         text_color=WHITE, size=10, letter_sp=500)
    add_text(s, Inches(2.25), Inches(1.55), Inches(4.5), Inches(0.42),
             'SHARE', size=11, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, letter_sp=600)

    add_text(s, Inches(0.55), Inches(2.2), Inches(6.4), Inches(0.9),
             'Share one',
             size=46, bold=True, color=WHITE)
    add_gradient_text(s, Inches(0.55), Inches(3.0), Inches(6.4), Inches(1.4),
                      'clean link.',
                      size=70, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=2700000)
    add_text(s, Inches(0.55), Inches(4.5), Inches(6.4), Inches(0.9),
             'Clients pick their favorites\nfrom any device.',
             size=18, color=MUTED, line_spacing=1.4)

    # featured link card with gradient outline
    lx, ly = Inches(0.6), Inches(6.1)
    lw, lh = Inches(6.3), Inches(1.4)
    add_rect(s, lx, ly, lw, lh, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, lx, ly, lw, lh,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=2.0, angle=2700000)
    icon_link(s, lx + Inches(0.3), ly + Inches(0.22), size_in=0.95,
              stops=[(0.0, ROYAL), (1.0, PURPLE)])
    add_text(s, lx + Inches(1.5), ly + Inches(0.22), lw - Inches(1.7), Inches(0.35),
             'Your album link',
             size=10, bold=True, color=MUTED, letter_sp=300)
    add_gradient_text(s, lx + Inches(1.5), ly + Inches(0.55), lw - Inches(1.7), Inches(0.7),
                      'framedrops.in/sneha-rohan',
                      size=22, bold=True,
                      stops=[(0.0, ROYAL), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=0)

    # device cards
    devices = [
        ('Mobile',  [(0.0, ROYAL), (1.0, PURPLE)]),
        ('Tablet',  [(0.0, PURPLE),(1.0, PINK)]),
        ('Desktop', [(0.0, PINK),  (1.0, ORANGE)]),
    ]
    dy = Inches(7.9)
    cw = Inches(2.0); ch = Inches(2.6); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    sx = (SW - total_w) // 2
    for i, (label, grad) in enumerate(devices):
        cx = sx + i * (cw + gap)
        add_rect(s, cx, dy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, cx, dy, cw, ch, grad, line_w_pt=1.5, angle=2700000)
        # tiny device mockup inside
        if label == 'Mobile':
            mw, mh = Inches(0.7), Inches(1.3)
        elif label == 'Tablet':
            mw, mh = Inches(1.1), Inches(1.3)
        else:
            mw, mh = Inches(1.5), Inches(1.0)
        mx = cx + (cw - mw) / 2
        my = dy + Inches(0.3)
        add_rect(s, mx, my, mw, mh, NAVY_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_rect(s, mx + Inches(0.06), my + Inches(0.08), mw - Inches(0.12), mh - Inches(0.16),
                          stops=grad, angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, cx, dy + ch - Inches(0.55), cw, Inches(0.4), label,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # tagline
    add_text(s, 0, Inches(10.8), SW, Inches(0.6),
             'No app install. No friction.',
             size=15, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, line_spacing=1.2)
    gradient_divider(s, Inches(11.6), w_in=2.0)


# ---------- Slide 6: No mess ----------
def slide_6_nomess(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.45), Inches(2.0), Inches(2.6), Inches(0.45),
         'GOODBYE CHAOS',
         gradient_stops=[(0.0, PINK), (1.0, ORANGE)],
         text_color=WHITE, size=11, letter_sp=400)

    add_gradient_text(s, Inches(0.4), Inches(3.0), Inches(6.7), Inches(2.0),
                      'No messy\nscreenshots.',
                      size=58, bold=True,
                      stops=[(0.0, PINK), (0.5, PURPLE), (1.0, ROYAL)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.05)
    add_text(s, Inches(0.4), Inches(5.3), Inches(6.7), Inches(0.9),
             'No endless chats.',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Before / After card
    bw = Inches(6.3); bh = Inches(5.0)
    bx = (SW - bw) // 2
    by = Inches(6.8)
    add_rect(s, bx, by, bw, bh, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, bx, by, bw, bh,
                           stops=[(0.0, ROYAL), (1.0, PINK)], line_w_pt=1.5, angle=2700000)

    # split header
    add_text(s, bx, by + Inches(0.25), bw / 2, Inches(0.4),
             'BEFORE', size=12, bold=True, color=PINK, align=PP_ALIGN.CENTER, letter_sp=500)
    add_text(s, bx + bw / 2, by + Inches(0.25), bw / 2, Inches(0.4),
             'WITH FRAMEDROPS', size=12, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, letter_sp=300)

    # divider
    add_rect(s, bx + bw / 2 - Inches(0.01), by + Inches(0.8), Inches(0.02), bh - Inches(1.1), DIVIDER)

    before = ['“Send pic 87 again”', 'WhatsApp clutter', 'Lost screenshots', 'Confused selections', 'Endless follow-ups']
    after  = ['One clean link', 'Just heart it ♥', 'Auto-saved picks', 'Crystal clear list', 'Done in one tap']

    for i in range(5):
        ry = by + Inches(1.1) + i * Inches(0.72)
        # before
        add_gradient_rect(s, bx + Inches(0.25), ry + Inches(0.18), Inches(0.18), Inches(0.18),
                          stops=[(0.0, PINK), (1.0, ORANGE)], angle=0, shape=MSO_SHAPE.OVAL)
        add_text(s, bx + Inches(0.55), ry, bw / 2 - Inches(0.7), Inches(0.55),
                 before[i], size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        # after
        add_gradient_rect(s, bx + bw / 2 + Inches(0.25), ry + Inches(0.18), Inches(0.18), Inches(0.18),
                          stops=[(0.0, ROYAL), (1.0, PURPLE)], angle=0, shape=MSO_SHAPE.OVAL)
        add_text(s, bx + bw / 2 + Inches(0.55), ry, bw / 2 - Inches(0.7), Inches(0.55),
                 after[i], size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 7: Heart ----------
def slide_7_heart(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.65), Inches(1.9), Inches(2.2), Inches(0.45),
         'CLIENT EXPERIENCE',
         gradient_stops=[(0.0, PINK), (1.0, PURPLE)],
         text_color=WHITE, size=11, letter_sp=400)

    add_text(s, Inches(0.55), Inches(2.8), Inches(6.4), Inches(0.9),
             'Clients simply',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, Inches(0.4), Inches(3.7), Inches(6.7), Inches(2.0),
                      'heart',
                      size=110, bold=True,
                      stops=[(0.0, PINK), (0.5, PURPLE), (1.0, ROYAL)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    # big heart icon
    hcx = SW / 2 - Inches(1.0)
    hcy = Inches(5.8)
    add_gradient_rect(s, hcx, hcy, Inches(2.0), Inches(2.0),
                      stops=[(0.0, PINK), (0.5, PURPLE), (1.0, ROYAL)],
                      angle=2700000, shape=MSO_SHAPE.OVAL)
    hs = s.shapes.add_shape(MSO_SHAPE.HEART,
                            hcx + Inches(0.3), hcy + Inches(0.35),
                            Inches(1.4), Inches(1.25))
    hs.fill.solid(); hs.fill.fore_color.rgb = WHITE
    hs.line.fill.background(); hs.shadow.inherit = False

    add_text(s, 0, Inches(8.1), SW, Inches(0.9),
             'the photos they want.',
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # phone mockup with gallery & hearts
    phone_w = Inches(2.6); phone_h = Inches(3.2)
    phone_x = (SW - phone_w) // 2
    phone_y = Inches(9.4)
    phone_mockup(s, phone_x, phone_y, phone_w, phone_h)
    # inside grid w/ hearts
    sx = phone_x + Inches(0.12); sy = phone_y + Inches(0.38)
    sw = phone_w - Inches(0.24); sh = phone_h - Inches(0.5)
    cols = 3; rows = 4; gap = Inches(0.05); pad = Inches(0.1)
    cw = (sw - 2 * pad - (cols - 1) * gap) / cols
    ch = cw
    palettes = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
        [(0.0, ROYAL), (1.0, PINK)],
        [(0.0, ORANGE), (1.0, PURPLE)],
        [(0.0, ROYAL), (1.0, ORANGE)],
    ]
    hearted = {(0, 1), (1, 0), (1, 2), (2, 1), (3, 0), (3, 2)}
    for r in range(rows):
        for c in range(cols):
            ccx = sx + pad + c * (cw + gap)
            ccy = sy + pad + r * (ch + gap)
            add_gradient_rect(s, ccx, ccy, cw, ch, palettes[(r * cols + c) % len(palettes)],
                              angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            if (r, c) in hearted:
                hh = s.shapes.add_shape(MSO_SHAPE.HEART,
                                        ccx + cw - Inches(0.22),
                                        ccy + ch - Inches(0.22),
                                        Inches(0.18), Inches(0.16))
                hh.fill.solid(); hh.fill.fore_color.rgb = WHITE
                hh.line.fill.background(); hh.shadow.inherit = False


# ---------- Slide 8: Local transfer (PRODUCT INVARIANT) ----------
def slide_8_transfer(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.45), Inches(1.7), Inches(2.6), Inches(0.45),
         'PRO MAGIC ✦',
         gradient_stops=[(0.0, ORANGE), (1.0, PINK)],
         text_color=WHITE, size=11, letter_sp=400)

    add_text(s, Inches(0.55), Inches(2.5), Inches(6.4), Inches(0.85),
             'Then Framedrops',
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, Inches(0.4), Inches(3.25), Inches(6.7), Inches(2.6),
                      'transfers your\noriginals locally',
                      size=44, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.05)
    add_text(s, Inches(0.55), Inches(5.8), Inches(6.4), Inches(0.9),
             'straight into your folders.',
             size=20, color=MUTED, align=PP_ALIGN.CENTER)

    # diagram: ALBUM → ARROW → LOCAL FOLDER
    dy = Inches(7.0)
    box_w = Inches(2.8); box_h = Inches(2.6)
    left_x = Inches(0.55)
    right_x = SW - box_w - Inches(0.55)

    # LEFT: Selected album card
    add_rect(s, left_x, dy, box_w, box_h, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, left_x, dy, box_w, box_h,
                           stops=[(0.0, PINK), (1.0, PURPLE)], line_w_pt=1.5, angle=2700000)
    add_text(s, left_x, dy + Inches(0.18), box_w, Inches(0.3),
             'SELECTED', size=9, bold=True, color=PINK, align=PP_ALIGN.CENTER, letter_sp=400)
    # mini grid w/ hearts
    gsx = left_x + Inches(0.25); gsy = dy + Inches(0.55)
    gsw = box_w - Inches(0.5)
    cols2 = 3; rows2 = 3
    cellw = (gsw - 2 * Inches(0.06)) / cols2
    for r in range(rows2):
        for c in range(cols2):
            ccx = gsx + c * (cellw + Inches(0.06))
            ccy = gsy + r * (cellw + Inches(0.06))
            pal = [[(0.0, ROYAL), (1.0, PURPLE)],
                   [(0.0, PURPLE), (1.0, PINK)],
                   [(0.0, PINK), (1.0, ORANGE)]][(r * cols2 + c) % 3]
            add_gradient_rect(s, ccx, ccy, cellw, cellw, pal, angle=2700000,
                              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            # tiny heart on every cell (all selected)
            hh = s.shapes.add_shape(MSO_SHAPE.HEART,
                                    ccx + cellw - Inches(0.16),
                                    ccy + cellw - Inches(0.16),
                                    Inches(0.13), Inches(0.11))
            hh.fill.solid(); hh.fill.fore_color.rgb = WHITE
            hh.line.fill.background(); hh.shadow.inherit = False

    # CENTER arrow + label
    ax = left_x + box_w
    aw = right_x - ax
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               ax + aw / 2 - Inches(0.7), dy + box_h / 2 - Inches(0.35),
                               Inches(1.4), Inches(0.7))
    # gradient fill arrow
    sp = arrow.fill._xPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, rgb in [(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', '0'); lin.set('scaled', '0')
    arrow.line.fill.background(); arrow.shadow.inherit = False
    add_text(s, ax, dy + box_h / 2 + Inches(0.45), aw, Inches(0.35),
             'on your device',
             size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=200)

    # RIGHT: Folder
    add_rect(s, right_x, dy, box_w, box_h, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, right_x, dy, box_w, box_h,
                           stops=[(0.0, ORANGE), (1.0, ROYAL)], line_w_pt=1.5, angle=2700000)
    add_text(s, right_x, dy + Inches(0.18), box_w, Inches(0.3),
             'YOUR FOLDER', size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, letter_sp=400)
    # big folder
    fcx = right_x + box_w / 2 - Inches(0.9); fcy = dy + Inches(0.7)
    add_gradient_rect(s, fcx, fcy, Inches(0.65), Inches(0.18),
                      stops=[(0.0, ORANGE), (1.0, PINK)], angle=2700000,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_rect(s, fcx - Inches(0.05), fcy + Inches(0.13), Inches(1.8), Inches(1.3),
                      stops=[(0.0, ORANGE), (1.0, PINK)], angle=2700000,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # filename
    add_text(s, right_x + Inches(0.15), dy + box_h - Inches(0.7), box_w - Inches(0.3), Inches(0.3),
             '/Sneha-Rohan/Selected',
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, right_x + Inches(0.15), dy + box_h - Inches(0.4), box_w - Inches(0.3), Inches(0.25),
             '124 originals · ready',
             size=8, color=MUTED, align=PP_ALIGN.CENTER)

    # footnote: invariant
    fy = Inches(10.0)
    add_rect(s, Inches(0.55), fy, Inches(6.4), Inches(0.65), CARD_FILL,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.55), fy, Inches(6.4), Inches(0.65),
             '⚡  Local transfer  •  Originals never re-downloaded',
             size=11, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 9: Payments ----------
def slide_9_payments(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(0.55), Inches(1.55), Inches(1.6), Inches(0.42),
         'STEP 04',
         gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)],
         text_color=WHITE, size=10, letter_sp=500)
    add_text(s, Inches(2.25), Inches(1.55), Inches(4.5), Inches(0.42),
             'GET PAID', size=11, bold=True, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE, letter_sp=600)

    add_text(s, Inches(0.55), Inches(2.2), Inches(6.4), Inches(0.9),
             'Collect',
             size=42, bold=True, color=WHITE)
    add_gradient_text(s, Inches(0.55), Inches(3.0), Inches(6.4), Inches(1.5),
                      'payments.',
                      size=72, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=2700000)
    add_text(s, Inches(0.55), Inches(4.7), Inches(6.4), Inches(0.6),
             'Right inside the album.',
             size=18, color=MUTED)

    # invoice hero card
    cx, cy = Inches(0.6), Inches(5.7)
    cw, ch = Inches(6.3), Inches(2.6)
    add_gradient_rect(s, cx, cy, cw, ch,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=2700000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, cx + Inches(0.4), cy + Inches(0.3), cw - Inches(0.8), Inches(0.4),
             'INVOICE  ·  SNEHA & ROHAN',
             size=10, bold=True, color=SOFT_WHITE, letter_sp=400)
    add_text(s, cx + Inches(0.4), cy + Inches(0.75), cw - Inches(0.8), Inches(1.0),
             '₹ 42,000',
             size=58, bold=True, color=WHITE)
    add_text(s, cx + Inches(0.4), cy + Inches(1.85), cw - Inches(0.8), Inches(0.4),
             'Wedding · Final delivery',
             size=12, color=SOFT_WHITE)
    # PAID stamp
    pill(s, cx + cw - Inches(1.55), cy + Inches(0.32), Inches(1.2), Inches(0.45),
         'PAID  ✓', gradient_stops=[(0.0, WHITE), (1.0, SOFT_WHITE)],
         text_color=PURPLE, size=12)

    # payment method cards
    methods = [
        ('UPI',        '₹ 18,000', [(0.0, ROYAL),  (1.0, PURPLE)],  icon_payment),
        ('Razorpay',   '₹ 18,000', [(0.0, PURPLE), (1.0, PINK)],    icon_lightning),
        ('Bank',       '₹  6,000', [(0.0, PINK),   (1.0, ORANGE)],  icon_check),
    ]
    my = Inches(8.7)
    mw = Inches(2.0); mh = Inches(2.4); gap = Inches(0.15)
    total_w = 3 * mw + 2 * gap
    msx = (SW - total_w) // 2
    for i, (label, amt, grad, icon_fn) in enumerate(methods):
        mx = msx + i * (mw + gap)
        add_rect(s, mx, my, mw, mh, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, mx, my, mw, mh, grad, line_w_pt=1.5, angle=2700000)
        icon_fn(s, mx + (mw - Inches(0.95)) / 2, my + Inches(0.3), size_in=0.95, stops=grad)
        add_text(s, mx, my + Inches(1.35), mw, Inches(0.3),
                 label, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_gradient_text(s, mx, my + Inches(1.7), mw, Inches(0.55),
                          amt, size=16, bold=True, stops=grad,
                          align=PP_ALIGN.CENTER)

    # bottom flourish
    add_text(s, 0, Inches(11.5), SW, Inches(0.4),
             'Razorpay · UPI · Bank',
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=500)


# ---------- Slide 10: Deliver finals ----------
def slide_10_deliver(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(0.55), Inches(1.55), Inches(1.6), Inches(0.42),
         'STEP 05',
         gradient_stops=[(0.0, PURPLE), (1.0, PINK)],
         text_color=WHITE, size=10, letter_sp=500)
    add_text(s, Inches(2.25), Inches(1.55), Inches(4.5), Inches(0.42),
             'DELIVER', size=11, bold=True, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE, letter_sp=600)

    add_text(s, Inches(0.55), Inches(2.2), Inches(6.4), Inches(0.9),
             'Deliver',
             size=44, bold=True, color=WHITE)
    add_gradient_text(s, Inches(0.55), Inches(3.0), Inches(6.4), Inches(1.5),
                      'the finals.',
                      size=72, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.LEFT, angle=2700000)

    add_text(s, Inches(0.55), Inches(4.6), Inches(6.4), Inches(0.5),
             'Edited. Polished. Done.',
             size=18, color=MUTED)

    # phone mockup with "Done" screen
    phone_w = Inches(3.6); phone_h = Inches(5.6)
    phone_x = (SW - phone_w) // 2
    phone_y = Inches(5.6)
    phone_mockup(s, phone_x, phone_y, phone_w, phone_h, screen_content=screen_done)

    # bottom stat row
    sy = Inches(11.6)
    cw = Inches(2.0); ch = Inches(0.95); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    sx = (SW - total_w) // 2
    stats = [('124', 'Originals'), ('100%', 'Quality'), ('1 click', 'Sent')]
    grads = [
        [(0.0, ROYAL), (1.0, PURPLE)],
        [(0.0, PURPLE), (1.0, PINK)],
        [(0.0, PINK), (1.0, ORANGE)],
    ]
    for i, (val, label) in enumerate(stats):
        cx = sx + i * (cw + gap)
        add_gradient_line_rect(s, cx, sy, cw, ch, grads[i], line_w_pt=1.5, angle=2700000)
        add_gradient_text(s, cx, sy + Inches(0.1), cw, Inches(0.45), val,
                          size=18, bold=True, stops=grads[i],
                          align=PP_ALIGN.CENTER)
        add_text(s, cx, sy + Inches(0.55), cw, Inches(0.35), label,
                 size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=200)


# ---------- Slide 11: DONE big ----------
def slide_11_done(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.75), Inches(3.2), Inches(2.0), Inches(0.45),
         'PROJECT WRAPPED',
         gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)],
         text_color=WHITE, size=11, letter_sp=400)

    add_gradient_text(s, 0, Inches(4.5), SW, Inches(4.5),
                      'Done.',
                      size=240, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.0,
                      letter_sp=-100)

    # checklist
    items = ['Uploaded', 'Shared', 'Selected', 'Transferred', 'Paid']
    cy = Inches(9.6)
    cw = Inches(6.0); ch = Inches(3.0)
    cx = (SW - cw) // 2
    add_rect(s, cx, cy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, cx, cy, cw, ch,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=1.5, angle=2700000)
    for i, label in enumerate(items):
        ry = cy + Inches(0.3) + i * Inches(0.5)
        # check circle
        add_gradient_rect(s, cx + Inches(0.4), ry + Inches(0.08), Inches(0.3), Inches(0.3),
                          stops=[(0.0, ROYAL), (1.0, PINK)], angle=2700000, shape=MSO_SHAPE.OVAL)
        add_text(s, cx + Inches(0.4), ry + Inches(0.08), Inches(0.3), Inches(0.3),
                 '✓', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.85), ry, cw - Inches(1.0), Inches(0.45),
                 label, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 12: Less chaos / more shoots ----------
def slide_12_lesschaos(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.45), Inches(2.0), Inches(2.6), Inches(0.45),
         'THE OUTCOME',
         gradient_stops=[(0.0, ORANGE), (1.0, PINK)],
         text_color=WHITE, size=11, letter_sp=400)

    add_text(s, 0, Inches(3.0), SW, Inches(1.0),
             'Less',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(3.8), SW, Inches(1.8),
                      'chaos.',
                      size=110, bold=True,
                      stops=[(0.0, PINK), (0.5, ORANGE), (1.0, PURPLE)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    add_text(s, 0, Inches(5.9), SW, Inches(1.0),
             'More',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(6.7), SW, Inches(1.8),
                      'shoots.',
                      size=110, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    add_text(s, 0, Inches(8.8), SW, Inches(1.0),
             'More',
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_text(s, 0, Inches(9.6), SW, Inches(1.8),
                      'business.',
                      size=110, bold=True,
                      stops=[(0.0, PURPLE), (0.5, PINK), (1.0, ORANGE)],
                      align=PP_ALIGN.CENTER, angle=2700000)

    # divider
    gradient_divider(s, Inches(11.7), w_in=2.4)


# ---------- Slide 13: Built for photographers ----------
def slide_13_builtfor(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.55), Inches(2.0), Inches(2.4), Inches(0.45),
         'WHO IT’S FOR',
         gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)],
         text_color=WHITE, size=11, letter_sp=400)

    add_gradient_text(s, Inches(0.4), Inches(2.95), Inches(6.7), Inches(2.0),
                      'Framedrops.',
                      size=80, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, letter_sp=-30)

    add_text(s, Inches(0.55), Inches(4.7), Inches(6.4), Inches(1.2),
             'Built for photographers\nwho move fast.',
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.3)

    # feature grid
    features = [
        ('Wedding pros',     icon_camera, [(0.0, ROYAL), (1.0, PURPLE)]),
        ('Event teams',      icon_lightning, [(0.0, PURPLE), (1.0, PINK)]),
        ('Family shoots',    icon_heart, [(0.0, PINK), (1.0, ORANGE)]),
        ('Brand & product',  icon_album, [(0.0, ORANGE), (1.0, ROYAL)]),
        ('Solo creators',    icon_chat, [(0.0, ROYAL), (1.0, PINK)]),
        ('Studio crews',     icon_folder, [(0.0, PURPLE), (1.0, ORANGE)]),
    ]
    gy = Inches(6.6)
    cw = Inches(2.0); ch = Inches(2.0); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    gx0 = (SW - total_w) // 2
    for i, (label, icon_fn, grad) in enumerate(features):
        col = i % 3
        row = i // 3
        gx = gx0 + col * (cw + gap)
        gyy = gy + row * (ch + gap)
        add_rect(s, gx, gyy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, gx, gyy, cw, ch, grad, line_w_pt=1.5, angle=2700000)
        icon_fn(s, gx + (cw - Inches(0.85)) / 2, gyy + Inches(0.35), size_in=0.85, stops=grad)
        add_text(s, gx, gyy + Inches(1.35), cw, Inches(0.45),
                 label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # tagline strip
    sy = Inches(11.4)
    add_rect(s, Inches(0.55), sy, Inches(6.4), Inches(0.65), CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.55), sy, Inches(6.4), Inches(0.65),
             'Fast. Beautiful. Made in India.',
             size=12, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=300)


# ---------- Slide 14: CTA ----------
def slide_14_cta(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    pill(s, Inches(2.95), Inches(2.6), Inches(1.6), Inches(0.45),
         'START NOW',
         gradient_stops=[(0.0, ROYAL), (1.0, PURPLE)],
         text_color=WHITE, size=11, letter_sp=400)

    add_text(s, 0, Inches(3.6), SW, Inches(0.9),
             'Try Framedrops today',
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # giant URL
    add_gradient_text(s, 0, Inches(4.7), SW, Inches(2.6),
                      'framedrops\n.in',
                      size=92, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, line_spacing=1.05,
                      letter_sp=-40)

    # CTA button
    bw, bh = Inches(4.6), Inches(1.0)
    bx = (SW - bw) // 2
    by = Inches(8.0)
    add_gradient_rect(s, bx, by, bw, bh,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      angle=0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, bx, by, bw, bh,
             'Start free  →',
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_sp=200)

    # 3 mini USP cards
    usps = [
        ('Free to start',   [(0.0, ROYAL), (1.0, PURPLE)]),
        ('No card needed',  [(0.0, PURPLE), (1.0, PINK)]),
        ('Setup in 2 min',  [(0.0, PINK), (1.0, ORANGE)]),
    ]
    uy = Inches(9.5)
    cw = Inches(2.0); ch = Inches(1.2); gap = Inches(0.15)
    total_w = 3 * cw + 2 * gap
    sx = (SW - total_w) // 2
    for i, (label, grad) in enumerate(usps):
        cx = sx + i * (cw + gap)
        add_rect(s, cx, uy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_gradient_line_rect(s, cx, uy, cw, ch, grad, line_w_pt=1.25, angle=2700000)
        # dot
        add_gradient_rect(s, cx + (cw - Inches(0.2)) / 2, uy + Inches(0.22), Inches(0.2), Inches(0.2),
                          stops=grad, angle=2700000, shape=MSO_SHAPE.OVAL)
        add_text(s, cx, uy + Inches(0.55), cw, Inches(0.55),
                 label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # social handle / footer
    add_text(s, 0, Inches(11.2), SW, Inches(0.4),
             '@framedrops  ·  framedrops.in',
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=400)

    # subtle bottom blob underline
    gradient_divider(s, Inches(11.85), w_in=2.4)


# ---------- Slide 15: Outro / signature ----------
def slide_15_outro(idx, total):
    s = prs.slides.add_slide(BLANK)
    background(s, 'dark')
    slide_chrome(s, idx, total)

    # big logo lockup
    add_gradient_text(s, 0, Inches(5.0), SW, Inches(2.0),
                      'Framedrops',
                      size=78, bold=True,
                      stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=2700000, letter_sp=-20)

    gradient_divider(s, Inches(7.2), w_in=2.2)

    add_text(s, 0, Inches(7.7), SW, Inches(0.6),
             'For photographers who move fast.',
             size=15, color=SOFT_WHITE, align=PP_ALIGN.CENTER, line_spacing=1.3)

    # follow card
    cw = Inches(5.5); ch = Inches(1.3)
    cx = (SW - cw) // 2
    cy = Inches(9.0)
    add_rect(s, cx, cy, cw, ch, CARD_FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_gradient_line_rect(s, cx, cy, cw, ch,
                           stops=[(0.0, ROYAL), (0.5, PURPLE), (1.0, PINK)],
                           line_w_pt=1.5, angle=2700000)
    add_text(s, cx, cy + Inches(0.2), cw, Inches(0.4),
             'FOLLOW', size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=500)
    add_gradient_text(s, cx, cy + Inches(0.55), cw, Inches(0.7),
                      '@framedrops',
                      size=26, bold=True,
                      stops=[(0.0, ROYAL), (1.0, PINK)],
                      align=PP_ALIGN.CENTER, angle=0)

    add_text(s, 0, Inches(11.0), SW, Inches(0.4),
             'framedrops.in',
             size=13, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, letter_sp=400)
    add_text(s, 0, Inches(11.5), SW, Inches(0.35),
             'Made with ♥ for photographers',
             size=10, color=MUTED, align=PP_ALIGN.CENTER, letter_sp=300)


# ---------- Build all slides ----------
TOTAL = 15
slide_1_hook(1, TOTAL)
slide_2_meet(2, TOTAL)
slide_3_upload(3, TOTAL)
slide_4_albums(4, TOTAL)
slide_5_share(5, TOTAL)
slide_6_nomess(6, TOTAL)
slide_7_heart(7, TOTAL)
slide_8_transfer(8, TOTAL)
slide_9_payments(9, TOTAL)
slide_10_deliver(10, TOTAL)
slide_11_done(11, TOTAL)
slide_12_lesschaos(12, TOTAL)
slide_13_builtfor(13, TOTAL)
slide_14_cta(14, TOTAL)
slide_15_outro(15, TOTAL)

out_path = '/Users/apple/Desktop/PHOTOSHARE/framedropnotes/Framedrops_Portrait_Deck.pptx'
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Slide size: {prs.slide_width} x {prs.slide_height} EMU '
      f'({prs.slide_width / 914400:.3f}" x {prs.slide_height / 914400:.3f}")')
print(f'Aspect: {prs.slide_width / prs.slide_height:.4f} (target 9/16 = {9/16:.4f})')
print(f'Slides: {len(prs.slides)}')
